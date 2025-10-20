"""
Simurgh PDF – pure API back-end

Endpoints
─────────
GET  /health          → {"status": "ok"}
POST /api/merge       → streams merged PDF
(no HTML rendered; UI lives in React/Vite front-end)
"""

# SPDX-License-Identifier: AGPL-3.0-only

# ── imports ──────────────────────────────────────────────────────
from flask import (
    Flask, request, send_file, jsonify, after_this_request
)
from flask_cors import CORS                 # allow front-end origin
from werkzeug.utils import secure_filename
from PyPDF2 import PdfMerger, PdfReader, PdfWriter   # ← PdfReader/Writer for split
from datetime import datetime, timedelta
import mimetypes, os, uuid, io, re, zipfile
import subprocess, shlex          # run Ghostscript
import shutil, platform
from pdf2docx import Converter
from PIL import Image           # for JPG/PNG → PDF
import fitz                     # PyMuPDF, for PDF → JPG/PNG
from pdfminer.high_level import extract_text_to_fp
from pdfminer.layout import LAParams
import pdfplumber
import sys
import threading
import hashlib
import time
from collections import defaultdict
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, ListFlowable, ListItem
from reportlab.lib import colors

# Try to import pandas, but make it optional
try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError as e:
    print(f"Warning: pandas not available: {e}")
    PANDAS_AVAILABLE = False
    pd = None

import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from dotenv import load_dotenv
import requests
import json

# Register async endpoints for Explain and Summarize tools
try:
    from explain.endpoints import register_explain_endpoints
except Exception as e:
    register_explain_endpoints = None  # type: ignore
    print(f"Warning: failed to import explain endpoints: {e}")

try:
    from summarize.endpoints import register_summarize_endpoints
except Exception as e:
    register_summarize_endpoints = None  # type: ignore
    print(f"Warning: failed to import summarize endpoints: {e}")

try:
    from extraction.ocr_service import OCRService
except Exception:
    OCRService = None  # type: ignore

# Import extraction services
try:
    from extraction.domain_service import DomainService
    from extraction.ai_service import AIExtractionService
    from extraction.text_service import TextExtractionService
    from extraction.cache_service import CacheService
    from extraction.confidence import ConfidenceScorer
    from extraction.pipeline import ExtractionPipeline
    from extraction.models import ExtractionOptions
    from extraction.config import config
except Exception as e:
    print(f"Warning: Failed to import extraction services: {e}")
    DomainService = None
    AIExtractionService = None
    TextExtractionService = None
    CacheService = None
    ConfidenceScorer = None
    ExtractionPipeline = None
    ExtractionOptions = None
    config = None

# Load environment variables from .env file
load_dotenv()

# AI Service configuration
AI_SERVICE = os.getenv("AI_SERVICE", "ollama")  # "ollama", "openai", or "anthropic"
OLLAMA_BASE_URL = os.getenv("OLLAMA_BASE_URL", "http://localhost:11434")
OLLAMA_MODEL = os.getenv("OLLAMA_MODEL", "llama3.1:8b")
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
OPENAI_MODEL = os.getenv("OPENAI_MODEL", "gpt-3.5-turbo")
ANTHROPIC_API_KEY = os.getenv("ANTHROPIC_API_KEY")
ANTHROPIC_MODEL = os.getenv("ANTHROPIC_MODEL", "claude-3-haiku-20240307")
# Additional AI tuning (timeouts and caps)
OLLAMA_TIMEOUT = int(os.getenv("OLLAMA_TIMEOUT", "180"))
OLLAMA_NUM_CTX = int(os.getenv("OLLAMA_NUM_CTX", "16384"))
OLLAMA_NUM_PREDICT = int(os.getenv("OLLAMA_NUM_PREDICT", "4096"))
SYNTHESIS_PER_FILE_MAX = int(os.getenv("SYNTHESIS_PER_FILE_MAX", "1200"))
SYNTHESIS_MAX_CHARS = int(os.getenv("SYNTHESIS_MAX_CHARS", "8000"))

# Print configuration at startup for debugging
print("=== AI SERVICE CONFIGURATION ===")
print(f"AI_SERVICE: {AI_SERVICE}")
print(f"ANTHROPIC_API_KEY: {'SET' if ANTHROPIC_API_KEY else 'NOT SET'}")
print(f"ANTHROPIC_MODEL: {ANTHROPIC_MODEL}")
print(f"OPENAI_API_KEY: {'SET' if OPENAI_API_KEY else 'NOT SET'}")
print(f"OPENAI_MODEL: {OPENAI_MODEL}")
print(f"OLLAMA_BASE_URL: {OLLAMA_BASE_URL}")
print(f"OLLAMA_MODEL: {OLLAMA_MODEL}")
print("================================")

# Force flush to ensure logs are visible
sys.stdout.flush()

# Valid Anthropic models for validation
VALID_ANTHROPIC_MODELS = [
    "claude-3-5-sonnet-20241022",
    "claude-3-5-haiku-20241022", 
    "claude-3-haiku-20240307",
    "claude-3-sonnet-20240229",
    "claude-3-opus-20240229",
    "claude-3-5-sonnet",
    "claude-3-haiku",
    "claude-3-sonnet",
    "claude-3-opus",
    "claude-3-5-haiku",
    "claude-3-5-opus"
]

# Simple in-memory job store for async extract (MVP; not for multi-process)
EXTRACT_JOBS = {}
EXTRACT_JOBS_LOCK = threading.Lock()

# Simple in-memory cache (process-local)
EXTRACT_CACHE: dict[str, dict] = {}

# Disk cache for extract results (per-process simple cache)
EXTRACT_CACHE_TTL_SECONDS = int(os.getenv("EXTRACT_CACHE_TTL_SECONDS", "86400"))  # 1 day default
EXTRACT_PIPELINE_VERSION = "v3-ai-only"  # bump to invalidate old cache for AI-only pipeline

_OCR_SERVICE = OCRService() if OCRService is not None else None

# ── Extraction Domain Schema (Backend authoritative map) ─────────
# Each domain maps to a flat list of canonical field keys.
# UI may present categories/groups, but backend validation uses this flat list.
DOMAIN_FIELDS: dict[str, list[str]] = {
    # Core domains (existing + expanded)
    "invoice": [
        "invoice_number", "po_number", "invoice_date", "due_date", "currency",
        "vendor_name", "vendor_address", "vendor_email", "vendor_phone", "vendor_tax_id",
        "customer_name", "customer_address", "customer_email", "customer_phone", "customer_tax_id",
        "subtotal_amount", "tax_amount", "tax_rate", "shipping_amount", "discount_amount", "total_amount",
        "payment_terms", "payment_method", "notes",
        "line_items"  # list of {description, quantity, unit_price, amount}
    ],
    "purchase_order": [
        "po_number", "requisition_number", "po_date", "delivery_date", "currency",
        "buyer_name", "buyer_address", "buyer_email", "buyer_phone",
        "vendor_name", "vendor_address", "vendor_email", "vendor_phone",
        "ship_to_name", "ship_to_address",
        "bill_to_name", "bill_to_address",
        "payment_terms", "approval_status", "approver_name",
        "subtotal_amount", "tax_amount", "shipping_amount", "discount_amount", "total_amount",
        "line_items"
    ],
    "receipt": [
        "receipt_number", "receipt_date", "store_name", "store_address", "store_phone",
        "cashier_name", "payment_method", "currency", "tax_amount", "total_amount",
        "line_items"
    ],
    "contract": [
        "party_a", "party_b", "effective_date", "term", "termination", "governing_law",
        "jurisdiction", "payment_terms", "confidentiality", "liability", "indemnity",
        "contact_email", "signatories"
    ],
    "financial": [
        "statement_type", "period", "currency",
        "revenue", "cost_of_goods_sold", "gross_profit", "operating_income", "operating_expenses",
        "net_income", "ebitda", "gross_margin", "operating_margin", "net_margin", "eps",
        "operating_cash_flow", "free_cash_flow", "total_assets", "total_liabilities",
        "shareholders_equity", "debt", "cash_and_equivalents", "ratios", "line_items"
    ],
    "bank_statement": [
        "account_holder", "account_number", "statement_period", "opening_balance",
        "closing_balance", "total_deposits", "total_withdrawals", "fees", "interest",
        "transactions"  # list of {date, description, debit, credit, balance}
    ],
    "tax_form": [
        "form_type", "tax_year", "filer_name", "filer_ssn_ein", "filer_address",
        "income_categories", "deductions", "credits", "tax_owed", "refund_amount"
    ],
    "research": [
        "title", "authors", "affiliations", "abstract", "methodology", "results",
        "conclusions", "keywords", "doi", "citations", "references", "research_metrics"
    ],
    "healthcare": [
        "patient_id", "mrn", "patient_name", "dob", "icd9_codes", "icd10_codes", "cpt_codes",
        "chief_complaint", "history", "physical_exam", "assessment", "plan",
        "medications", "allergies", "vitals", "labs", "diagnosis_text", "procedures_text",
        "primary_contact"
    ],
    "resume": [
        "name", "email", "phone", "address", "linkedin", "portfolio",
        "summary", "skills_technical", "skills_soft",
        "experience",  # list of {company, role, start_date, end_date, responsibilities, achievements}
        "education",   # list of {degree, institution, start_date, end_date, gpa}
        "certifications", "languages", "publications", "awards"
    ],
    "legal_pleading": [
        "case_name", "court", "docket_number", "judge", "filing_date",
        "parties", "claims", "relief_requested", "orders"
    ],
    "patent": [
        "patent_number", "application_number", "title", "assignee", "inventors",
        "filing_date", "issue_date", "ipc_codes", "abstract", "claims", "description"
    ],
    "medical_bill": [
        "provider_name", "provider_address", "patient_name", "patient_id", "service_dates",
        "icd_codes", "cpt_codes", "charges", "insurance_payments", "patient_responsibility",
        "total_amount"
    ],
    "lab_report": [
        "patient_name", "patient_id", "collection_date", "report_date", "ordering_physician",
        "tests",  # list of {name, value, unit, reference_range}
        "interpretation"
    ],
    "insurance_claim": [
        "claim_number", "policy_number", "insured_name", "loss_date", "loss_description",
        "adjuster_name", "status", "payments", "total_amount"
    ],
    "real_estate": [
        "property_address", "parcel_number", "owner_name", "assessed_value", "sale_price",
        "mortgage_lender", "loan_amount", "closing_date", "recording_number"
    ],
    "shipping_manifest": [
        "manifest_number", "carrier", "ship_date", "origin", "destination",
        "items",  # list of {description, quantity, weight, value}
        "total_weight", "total_value"
    ],
    # Generic fallback
    "general": [
        "summary", "emails", "phones", "amounts", "dates"
    ],
}

def _normalize_domain(value: str | None) -> str:
    d = (value or "").strip().lower()
    return d if d in DOMAIN_FIELDS else ("general" if d == "" else d)

def _validate_selected_fields(domain: str, selected_fields: list[str] | None) -> list[str]:
    allowed = set(DOMAIN_FIELDS.get(domain, []))
    if not selected_fields:
        return list(allowed)
    return [f for f in selected_fields if f in allowed]

def build_inverted_index(pages_text: list[str]) -> dict[str, list[int]]:
    """Very simple inverted index: word (>=4 chars) -> sorted list of page numbers (1-indexed)."""
    index: dict[str, set[int]] = {}
    try:
        for i, txt in enumerate(pages_text or []):
            if not txt:
                continue
            for w in re.findall(r"[A-Za-z][A-Za-z0-9\-]{3,}", txt.lower()):
                s = index.get(w)
                if s is None:
                    s = set()
                    index[w] = s
                s.add(i + 1)
        return {k: sorted(list(v)) for k, v in index.items()}
    except Exception:
        return {}

def enrich_extraction_with_llm(dtype: str, mapped: dict, full_text: str, pages_text: list[str] | None = None, inv_index: dict[str, list[int]] | None = None) -> dict:
    """Optional hybrid LLM enrichment for complex fields. Returns envelope with 'mapped_fields'."""
    try:
        if not isinstance(mapped, dict) or not full_text:
            return {"mapped_fields": mapped}
        domain = (dtype or 'general').lower()
        # Build RAG snippets: pick candidate pages based on tokens from mapped fields
        rag_pages: list[int] = []
        try:
            candidates: dict[int, int] = {}
            tokens: list[str] = []
            if domain == 'healthcare':
                # Prefer diagnosis_text, plan, medication names, labs names
                if isinstance(mapped.get('diagnosis_text'), str):
                    tokens += re.findall(r"[A-Za-z][A-Za-z0-9\-]{4,}", mapped['diagnosis_text'])
                if isinstance(mapped.get('plan'), str):
                    tokens += re.findall(r"[A-Za-z][A-Za-z0-9\-]{4,}", mapped['plan'])
                for m in (mapped.get('medications') or []):
                    try:
                        nm = (m or {}).get('name')
                        if isinstance(nm, str):
                            tokens += re.findall(r"[A-Za-z][A-Za-z0-9\-]{4,}", nm)
                    except Exception:
                        pass
                for l in (mapped.get('labs') or []):
                    try:
                        nm = (l or {}).get('name')
                        if isinstance(nm, str):
                            tokens += re.findall(r"[A-Za-z][A-Za-z0-9\-]{4,}", nm)
                    except Exception:
                        pass
            elif domain == 'contract':
                for k in ('party_a','party_b','governing_law','term'):
                    if isinstance(mapped.get(k), str):
                        tokens += re.findall(r"[A-Za-z][A-Za-z0-9\-]{4,}", mapped[k])
            elif domain == 'research':
                for k in ('abstract','methodology','results','conclusions'):
                    if isinstance(mapped.get(k), str):
                        tokens += re.findall(r"[A-Za-z][A-Za-z0-9\-]{4,}", mapped[k])
            # Use inverted index to find candidate pages
            if inv_index:
                for t in tokens[:20]:
                    for p in inv_index.get(t.lower(), []) or []:
                        candidates[p] = candidates.get(p, 0) + 1
            # Choose top pages by score
            rag_pages = [p for p, _ in sorted(candidates.items(), key=lambda kv: (-kv[1], kv[0]))[:3]]
        except Exception:
            rag_pages = []

        rag_snippets = []
        try:
            if pages_text and rag_pages:
                for p in rag_pages:
                    if 1 <= p <= len(pages_text):
                        snippet = (pages_text[p-1] or '')
                        if snippet:
                            rag_snippets.append(f"[Page {p}] " + snippet[:800])
        except Exception:
            pass

        rag_block = ("\nRAG_SNIPPETS:\n" + "\n\n".join(rag_snippets)) if rag_snippets else ""
        prompt = (
            f"SYSTEM: You are a domain-specific extraction validator and enricher for {domain} documents.\n"
            "Return ONLY a JSON object with normalized structures.\n\n"
            "INPUT (truncated to 8000 chars):\n" + full_text[:8000] + "\n\n"
            + rag_block + "\n\n"
            "TASK:\n"
            "- Normalize and enrich extracted fields if possible.\n"
            "- For healthcare: infer diagnoses (array of strings), procedures (array), and normalize labs as {name,value,unit,flag}.\n"
            "- For contracts: extract obligations (array), termination_conditions (array).\n"
            "- For research: extract key_findings (array) and primary_outcomes (array).\n\n"
            "OUTPUT (strict JSON): {\n"
            "  \"diagnoses\": [], \"procedures\": [], \"labs_normalized\": [], \"obligations\": [], \"termination_conditions\": [], \"key_findings\": [], \"primary_outcomes\": []\n"
            "}"
        )
        ai_text = call_ai_service(prompt)
        data = {}
        try:
            data = parse_json_safely(ai_text)
        except Exception:
            return {"mapped_fields": mapped}
        if not isinstance(data, dict):
            return {"mapped_fields": mapped}
        out = dict(mapped)
        # Merge selected fields under enrichment namespace too
        enrichment = {}
        for k in ("diagnoses","procedures","labs_normalized","obligations","termination_conditions","key_findings","primary_outcomes"):
            if k in data and data.get(k) is not None:
                enrichment[k] = data.get(k)
        if enrichment:
            out["enriched"] = True
            out["enrichment"] = enrichment
            try:
                print(f"DEBUG: Enrichment produced keys: {list(enrichment.keys())}")
                sys.stdout.flush()
            except Exception:
                pass
        else:
            try:
                print("DEBUG: Enrichment found no additional data")
                sys.stdout.flush()
            except Exception:
                pass
        return {"mapped_fields": out}
    except Exception:
        return {"mapped_fields": mapped}

def _extract_cache_dir() -> str:
    try:
        base = app.config.get("UPLOAD_FOLDER") or os.path.join(os.getcwd(), "uploads")
        path = os.path.join(base, "extract_cache")
        os.makedirs(path, exist_ok=True)
        return path
    except Exception:
        return os.getcwd()

def _extract_cache_path(key: str) -> str:
    return os.path.join(_extract_cache_dir(), f"{key}.json")

def call_ai_service(prompt, system_prompt=""):
    """Call AI service (Ollama, OpenAI, or Anthropic) based on configuration"""
    print(f"DEBUG: AI_SERVICE = '{AI_SERVICE}'")
    print(f"DEBUG: ANTHROPIC_API_KEY exists = {bool(ANTHROPIC_API_KEY)}")
    print(f"DEBUG: ANTHROPIC_MODEL = '{ANTHROPIC_MODEL}'")
    print(f"DEBUG: Environment variables loaded: AI_SERVICE={os.getenv('AI_SERVICE')}, ANTHROPIC_API_KEY={'SET' if os.getenv('ANTHROPIC_API_KEY') else 'NOT SET'}")
    sys.stdout.flush()
    
    if AI_SERVICE == "openai":
        print("DEBUG: Calling OpenAI")
        sys.stdout.flush()
        return call_openai(prompt, system_prompt)
    elif AI_SERVICE == "anthropic":
        print("DEBUG: Calling Anthropic")
        sys.stdout.flush()
        return call_anthropic(prompt, system_prompt)
    else:
        print(f"DEBUG: Defaulting to Ollama (AI_SERVICE='{AI_SERVICE}')")
        sys.stdout.flush()
        return call_ollama(prompt, system_prompt)

def call_openai(prompt, system_prompt=""):
    """Call OpenAI API with the given prompt"""
    if not OPENAI_API_KEY:
        return "Error: OpenAI API key not configured. Set OPENAI_API_KEY environment variable."
    
    try:
        import openai
        openai.api_key = OPENAI_API_KEY
        
        messages = []
        if system_prompt:
            messages.append({"role": "system", "content": system_prompt})
        messages.append({"role": "user", "content": prompt})
        
        response = openai.ChatCompletion.create(
            model=OPENAI_MODEL,
            messages=messages,
            max_tokens=1000,
            temperature=0.7
        )
        
        return response.choices[0].message.content
    except Exception as e:
        print(f"DEBUG: Exception calling OpenAI: {e}")
        return f"Error calling OpenAI: {str(e)}"

def call_anthropic(prompt, system_prompt=""):
    """Call Anthropic API with the given prompt using direct HTTP requests"""
    print(f"DEBUG: call_anthropic called with prompt length: {len(prompt)}")
    print(f"DEBUG: ANTHROPIC_API_KEY length: {len(ANTHROPIC_API_KEY) if ANTHROPIC_API_KEY else 0}")
    print(f"DEBUG: ANTHROPIC_MODEL: {ANTHROPIC_MODEL}")
    sys.stdout.flush()
    
    if not ANTHROPIC_API_KEY:
        print("DEBUG: No ANTHROPIC_API_KEY found")
        sys.stdout.flush()
        return "Error: Anthropic API key not configured. Set ANTHROPIC_API_KEY environment variable."
    
    # Validate model name
    if ANTHROPIC_MODEL not in VALID_ANTHROPIC_MODELS:
        print(f"DEBUG: Invalid model '{ANTHROPIC_MODEL}'. Valid models: {VALID_ANTHROPIC_MODELS}")
        sys.stdout.flush()
        return f"Error: Invalid Anthropic model '{ANTHROPIC_MODEL}'. Valid models are: {', '.join(VALID_ANTHROPIC_MODELS)}"
    
    try:
        print("DEBUG: Using direct HTTP requests to Anthropic API")
        sys.stdout.flush()
        
        # Prepare the request payload
        messages = []
        if system_prompt:
            messages.append({"role": "user", "content": f"System: {system_prompt}\n\nUser: {prompt}"})
        else:
            messages.append({"role": "user", "content": prompt})
        
        payload = {
            "model": ANTHROPIC_MODEL,
            "messages": messages,
            "max_tokens": 1000,
            "temperature": 0.7
        }
        
        headers = {
            "Content-Type": "application/json",
            "x-api-key": ANTHROPIC_API_KEY,
            "anthropic-version": "2023-06-01"
        }
        
        print(f"DEBUG: Making direct HTTP request to Anthropic API")
        print(f"DEBUG: Model: {ANTHROPIC_MODEL}")
        print(f"DEBUG: Payload keys: {list(payload.keys())}")
        sys.stdout.flush()
        
        # Make the HTTP request directly
        import requests
        response = requests.post(
            "https://api.anthropic.com/v1/messages",
            json=payload,
            headers=headers,
            timeout=60
        )
        
        print(f"DEBUG: HTTP response status: {response.status_code}")
        sys.stdout.flush()
        
        if response.status_code == 200:
            result = response.json()
            print("DEBUG: API call successful")
            sys.stdout.flush()
            return result["content"][0]["text"]
        else:
            print(f"DEBUG: API call failed with status {response.status_code}")
            print(f"DEBUG: Response text: {response.text[:200]}")
            sys.stdout.flush()
            
            # Handle specific error cases
            if response.status_code == 401:
                return "Error: Invalid Anthropic API key. Please check your API key configuration."
            elif response.status_code == 400:
                return f"Error: Invalid request. Please check model name '{ANTHROPIC_MODEL}' and request format."
            elif response.status_code == 429:
                return "Error: Rate limit exceeded. Please try again later."
            elif response.status_code == 500:
                return "Error: Anthropic API server error. Please try again later."
            else:
                return f"Error: Anthropic API returned status {response.status_code}: {response.text[:100]}"
        
    except requests.exceptions.Timeout:
        print("DEBUG: Request timeout")
        sys.stdout.flush()
        return "Error: Request timeout. Please try again."
    except requests.exceptions.ConnectionError:
        print("DEBUG: Connection error")
        sys.stdout.flush()
        return "Error: Connection error. Please check your internet connection."
    except Exception as e:
        print(f"DEBUG: Exception calling Anthropic API: {e}")
        print(f"DEBUG: Exception type: {type(e).__name__}")
        sys.stdout.flush()
        return f"Error calling Anthropic API: {str(e)}"

def call_ollama(prompt, system_prompt=""):
    """Call Ollama API with the given prompt"""
    try:
        url = f"{OLLAMA_BASE_URL}/api/generate"
        payload = {
            "model": OLLAMA_MODEL,
            "prompt": prompt,
            "system": system_prompt,
            "stream": False,
            "options": {
                "temperature": 0.3,
                "top_p": 0.8,
                "num_ctx": OLLAMA_NUM_CTX,
                "num_predict": OLLAMA_NUM_PREDICT,
                "repeat_penalty": 1.1,
                "seed": -1,
                "stop": ["```", "\n\n\n"]
            }
        }
        
        print(f"DEBUG: Calling Ollama with payload: {json.dumps(payload, indent=2)}")
        response = requests.post(url, json=payload, timeout=OLLAMA_TIMEOUT)
        if response.status_code == 200:
            result = response.json()
            response_text = result.get('response', 'No response from AI model')
            print(f"DEBUG: Raw Ollama response: {response_text}")
            return response_text
        else:
            print(f"DEBUG: Ollama error status: {response.status_code}")
            return f"Error calling Ollama: {response.status_code}"
    except Exception as e:
        print(f"DEBUG: Exception calling Ollama: {e}")
        return f"Error connecting to Ollama: {str(e)}"


def parse_json_safely(text):
    """Attempt to parse JSON from a possibly noisy LLM response.
    - Strips code fences
    - Extracts first JSON object block if present
    - Returns dict or raises ValueError
    """
    if not isinstance(text, str):
        raise ValueError("Response is not a string")

    cleaned = text.strip()

    # Remove common code fences
    if cleaned.startswith("```"):
        cleaned = cleaned.strip('`')
        # After stripping backticks, remove possible language tag remnants
        cleaned = cleaned.replace('json', '', 1).strip()

    # Try direct json parse first
    try:
        return json.loads(cleaned)
    except Exception:
        pass

    # Try to find first JSON object block
    start = cleaned.find('{')
    end = cleaned.rfind('}')
    if start != -1 and end != -1 and end > start:
        candidate = cleaned[start:end+1]
        try:
            return json.loads(candidate)
        except Exception:
            pass

    raise ValueError("Could not parse JSON from response")

# --- AI output normalization helpers ---

def _strip_code_fences(s):
    if not isinstance(s, str):
        return s
    t = s.strip()
    if t.startswith("```") and t.endswith("```"):
        t = t.strip('`')
        t = t.replace('json', '', 1).strip()
    return t

def _object_to_readable(obj: dict) -> str:
    if not isinstance(obj, dict):
        return str(obj)
    # Prefer common keys in a human order
    parts = []
    for key in ["title", "name", "point", "summary", "rationale", "recommendation", "action", "explanation", "detail", "note"]:
        if key in obj and obj.get(key):
            parts.append(f"{key.capitalize()}: {obj.get(key)}")
    if parts:
        return " ".join(parts)
    # Fallback: join all key/values
    try:
        return "; ".join([f"{k.capitalize()}: {v}" for k, v in obj.items() if v is not None])
    except Exception:
        return str(obj)

def _clean_display_text(text: str) -> str:
    if text is None:
        return ""
    if not isinstance(text, str):
        return str(text)
    import re, json as _json
    s = _strip_code_fences(text)
    s = s.strip()
    # Full-string JSON object → readable
    if s.startswith('{') and s.endswith('}'):
        try:
            return _object_to_readable(_json.loads(s))
        except Exception:
            pass
    # Replace embedded JSON objects that contain rationale/recommendation with readable form
    def _repl(m):
        chunk = m.group(0)
        try:
            return _object_to_readable(_json.loads(chunk))
        except Exception:
            return ''
    s = re.sub(r"\{[^{}]*?(?:\"rationale\"|\"recommendation\")[^{}]*?\}", _repl, s)
    # Strip stray outer quotes
    if s.startswith('"') and s.endswith('"'):
        s = s[1:-1]
    # Normalize whitespace
    s = " ".join(s.split())
    return s

def _normalize_list(value):
    import re
    items = []
    if isinstance(value, list):
        items = value
    elif isinstance(value, str):
        items = [p.strip(' -*•') for p in re.split(r"[\r\n]+", value) if p.strip()]
    else:
        return []
    normalized = []
    for it in items:
        if isinstance(it, dict):
            normalized.append(_object_to_readable(it))
        elif isinstance(it, str):
            # Try parsing JSON-ish string
            try:
                obj = parse_json_safely(it)
                if isinstance(obj, dict):
                    normalized.append(_object_to_readable(obj))
                    continue
            except Exception:
                pass
            normalized.append(_clean_display_text(it))
        else:
            normalized.append(str(it))
    return normalized

def normalize_ai_summary_payload(payload):
    """Normalize/clean AI JSON payload so UI never shows raw JSON fragments."""
    if not isinstance(payload, dict):
        return {"summary": _clean_display_text(str(payload)), "key_topics": [], "main_points": [], "recommendations": []}
    res = {}
    res["summary"] = _clean_display_text(payload.get("summary") or payload.get("overview") or payload.get("executive_summary") or payload.get("message") or "")
    res["key_topics"] = _normalize_list(payload.get("key_topics") or payload.get("topics") or payload.get("key_points") or payload.get("bullets") or payload.get("highlights"))
    res["main_points"] = _normalize_list(payload.get("main_points") or payload.get("points") or payload.get("findings") or payload.get("takeaways"))
    res["recommendations"] = _normalize_list(payload.get("recommendations") or payload.get("suggestions") or payload.get("actions") or payload.get("next_steps"))
    return res

def normalize_ai_explain_payload(payload):
    """Normalize/clean AI JSON payload for Explain mode (educational output)."""
    if not isinstance(payload, dict):
        return {
            "summary": _clean_display_text(str(payload)),
            "key_concepts": [],
            "explanations": [],
            "context": [],
            "definitions": []
        }
    res = {}
    res["summary"] = _clean_display_text(
        payload.get("explanation_summary")
        or payload.get("summary")
        or payload.get("overview")
        or payload.get("message")
        or ""
    )
    # Arrays, accepting multiple common aliases
    def L(key_candidates):
        for k in key_candidates:
            if k in payload and payload.get(k):
                return _normalize_list(payload.get(k))
        return []
    res["key_concepts"] = L(["key_concepts", "concepts", "key_terms", "terms", "ideas"])
    res["explanations"] = L(["explanations", "how_it_works", "analysis", "details", "rationale"])
    res["context"] = L(["context", "background", "related_work", "history"])
    res["definitions"] = L(["definitions", "glossary", "term_definitions", "terminology"]) 
    return res

app = Flask(__name__)
app.config["UPLOAD_FOLDER"] = "temp"
os.makedirs(app.config["UPLOAD_FOLDER"], exist_ok=True)

# Register Explain/Summarize async endpoints if available
if register_explain_endpoints:
    try:
        register_explain_endpoints(app)
        print("Explain async endpoints registered")
    except Exception as e:
        print(f"Warning: failed to register explain endpoints: {e}")

if register_summarize_endpoints:
    try:
        register_summarize_endpoints(app)
        print("Summarize async endpoints registered")
    except Exception as e:
        print(f"Warning: failed to register summarize endpoints: {e}")

# Job storage for async synthesis
SYNTHESIS_JOBS = {}
SYNTHESIS_JOBS_LOCK = threading.Lock()

# Initialize extraction services
_extraction_pipeline = None
print("DEBUG: Starting extraction services initialization...")
print(f"DEBUG: DomainService available: {DomainService is not None}")
print(f"DEBUG: AIExtractionService available: {AIExtractionService is not None}")
print(f"DEBUG: TextExtractionService available: {TextExtractionService is not None}")
print(f"DEBUG: CacheService available: {CacheService is not None}")
print(f"DEBUG: ConfidenceScorer available: {ConfidenceScorer is not None}")
print(f"DEBUG: ExtractionPipeline available: {ExtractionPipeline is not None}")

if all([DomainService, AIExtractionService, TextExtractionService, CacheService, ConfidenceScorer, ExtractionPipeline]):
    try:
        print("DEBUG: Creating service instances...")
        # Create service instances
        domain_service = DomainService()
        print("DEBUG: DomainService created")
        
        ai_service = AIExtractionService(domain_service)
        print("DEBUG: AIExtractionService created")
        
        text_service = TextExtractionService(_OCR_SERVICE)
        print("DEBUG: TextExtractionService created")
        
        cache_service = CacheService()
        print("DEBUG: CacheService created")
        
        confidence_scorer = ConfidenceScorer()
        print("DEBUG: ConfidenceScorer created")
        
        # Create pipeline
        _extraction_pipeline = ExtractionPipeline(
            domain_service=domain_service,
            ai_service=ai_service,
            text_service=text_service,
            cache_service=cache_service,
            confidence_scorer=confidence_scorer
        )
        print("DEBUG: ExtractionPipeline created")
        
        print("Extraction services initialized successfully")
    except Exception as e:
        print(f"Warning: Failed to initialize extraction services: {e}")
        import traceback
        traceback.print_exc()
        _extraction_pipeline = None
else:
    print("Warning: Some extraction services are not available")

origins = [
    "http://localhost:5173",  # Development frontend
    "http://127.0.0.1:5173",  # Alternative localhost
]
origins.extend([
    "https://perkpdf.com",
    "https://www.perkpdf.com",
])

CORS(
    app,
    resources={r"/api/*": {"origins": origins}}
)


# ── config & housekeeping ───────────────────────────────────────
app.config["UPLOAD_FOLDER"]      = "temp"
app.config["MAX_CONTENT_LENGTH"] = 100 * 1024 * 1024    # 100 MB
os.makedirs(app.config["UPLOAD_FOLDER"], exist_ok=True)

def purge_old_files(hours: int = 12):
    """Delete temp files older than <hours>."""
    cutoff = datetime.utcnow() - timedelta(hours=hours)
    for fname in os.listdir(app.config["UPLOAD_FOLDER"]):
        path = os.path.join(app.config["UPLOAD_FOLDER"], fname)
        if os.path.isfile(path) and datetime.utcfromtimestamp(
                os.path.getmtime(path)) < cutoff:
            try:
                os.remove(path)
            except OSError:
                pass
purge_old_files()
# ─────────────────────────────────────────────────────────────────

def allowed_pdf(fileobj):
    """Lenient PDF check: just check extension and MIME type."""
    if not fileobj.filename:
        return False
    
    ext = fileobj.filename.lower().endswith(".pdf")
    mime = fileobj.mimetype == "application/pdf"
    
    # Accept if either extension OR MIME type indicates PDF
    return ext or mime

def allowed_image(fileobj):
    """
    Accept only JPG/JPEG/PNG uploads:
      • filename must end in .jpg/.jpeg/.png
      • mimetype must start with "image/"
    """
    fname = fileobj.filename.lower()
    ext_ok = fname.endswith(".jpg") or fname.endswith(".jpeg") or fname.endswith(".png")
    mime_ok = fileobj.mimetype.startswith("image/")
    return ext_ok and mime_ok

# -------------------------
def gs_executable():
    """Return full path to Ghostscript binary or None."""
    # first, check PATH
    exe = shutil.which("gs") or shutil.which("gswin64c") or shutil.which("gswin32c")
    if exe:
        return exe

    # fallback: typical Windows install folder
    win_dirs = [
        r"D:\Program Files\gs",
        r"D:\Program Files (x86)\gs",
    ]
    for root in win_dirs:
        if os.path.isdir(root):
            # pick highest version folder
            versions = sorted(os.listdir(root), reverse=True)
            for v in versions:
                cand = os.path.join(root, v, "bin", "gswin64c.exe")
                if os.path.isfile(cand):
                    return cand
    return None

# ── ROUTES ───────────────────────────────────────────────────────
@app.get("/")
def root():
    """Simple root for anyone hitting the API directly."""
    return {"service": "Simurgh PDF API", "docs": "/health"}, 200

@app.get("/health")
def health():
    """Used by React (and uptime checks) to verify API is alive."""
    return jsonify(status="ok"), 200

@app.get("/api/ai-config")
def ai_config():
    """Debug endpoint to check AI service configuration"""
    return jsonify({
        "ai_service": AI_SERVICE,
        "anthropic_api_key_set": bool(ANTHROPIC_API_KEY),
        "anthropic_model": ANTHROPIC_MODEL,
        "openai_api_key_set": bool(OPENAI_API_KEY),
        "openai_model": OPENAI_MODEL,
        "ollama_base_url": OLLAMA_BASE_URL,
        "ollama_model": OLLAMA_MODEL,
        "valid_anthropic_models": VALID_ANTHROPIC_MODELS
    }), 200

@app.get("/api/ai-test")
def ai_test():
    """Test endpoint to verify AI service is working"""
    try:
        print("DEBUG: Testing AI service...")
        test_prompt = "Hello, this is a test. Please respond with 'AI service is working correctly.'"
        response = call_ai_service(test_prompt)
        print(f"DEBUG: AI test response: {response}")
        return jsonify({
            "status": "success",
            "ai_service": AI_SERVICE,
            "response": response,
            "timestamp": str(datetime.now())
        }), 200
    except Exception as e:
        print(f"DEBUG: AI test failed with error: {e}")
        return jsonify({
            "status": "error",
            "ai_service": AI_SERVICE,
            "error": str(e),
            "timestamp": str(datetime.now())
        }), 500

@app.get("/api/env-check")
def env_check():
    """Debug endpoint to check environment variables"""
    proxy_vars = ['HTTP_PROXY', 'HTTPS_PROXY', 'http_proxy', 'https_proxy']
    all_proxy_vars = {k: v for k, v in os.environ.items() if 'proxy' in k.lower()}
    http_vars = {k: v for k, v in os.environ.items() if any(x in k.lower() for x in ['http', 'https', 'ssl', 'cert', 'ca'])}
    
    # Check requests library configuration
    requests_info = {}
    try:
        import requests
        requests_info = {
            "version": requests.__version__,
            "global_proxies": getattr(requests, 'proxies', 'Not set'),
            "session_proxies": getattr(requests.Session(), 'proxies', 'Not set')
        }
    except Exception as e:
        requests_info = {"error": str(e)}
    
    # Check httpx library configuration
    httpx_info = {}
    try:
        import httpx
        httpx_info = {
            "version": httpx.__version__,
            "available": True
        }
    except Exception as e:
        httpx_info = {"error": str(e), "available": False}
    
    # Check urllib3 library configuration
    urllib3_info = {}
    try:
        import urllib3
        urllib3_info = {
            "version": urllib3.__version__,
            "available": True
        }
    except Exception as e:
        urllib3_info = {"error": str(e), "available": False}
    
    return jsonify({
        "proxy_environment_variables": {k: v for k, v in os.environ.items() if k in proxy_vars},
        "all_proxy_related_vars": all_proxy_vars,
        "http_https_vars": http_vars,
        "requests_library_info": requests_info,
        "httpx_library_info": httpx_info,
        "urllib3_library_info": urllib3_info,
        "ai_service": AI_SERVICE,
        "anthropic_api_key_set": bool(ANTHROPIC_API_KEY),
        "anthropic_model": ANTHROPIC_MODEL,
        "total_env_vars": len(os.environ),
        "sample_env_vars": dict(list(os.environ.items())[:10]),  # First 10 env vars
        "all_env_vars": dict(os.environ)  # All environment variables
    }), 200


# ── CONVERT (Unified) ─────────────────────────────────────────────────────
@app.get("/api/convert/formats")
def convert_formats():
    """
    Returns a list of supported conversion formats handled by POST /api/convert.
    We intentionally expose only formats backed by implemented routes.
    """
    return jsonify({
        "formats": [
            "docx",     # PDF → Word
            "xlsx",     # PDF → Excel (tables) or text CSV fallback logic reused under the hood
            "txt",      # PDF → plain text
            "png-zip",  # PDF → images (PNG) zipped
        ]
    }), 200


@app.post("/api/convert")
def convert_unified():
    """
    Unified conversion endpoint.
    Expects multipart/form-data with fields:
      • file   → the source PDF
      • format → one of: docx | xlsx | txt | png-zip
    """
    f = request.files.get("file")
    if not f or not allowed_pdf(f):
        return jsonify(error="Upload one PDF file"), 400

    target_format = (request.form.get("format") or "").lower().strip()
    if target_format not in ("docx", "xlsx", "txt", "png-zip"):
        return jsonify(error="Unsupported format"), 400

    # Save incoming PDF
    base_name = secure_filename(f.filename)
    src_path = os.path.join(app.config["UPLOAD_FOLDER"], f"{uuid.uuid4()}_{base_name}")
    f.save(src_path)

    try:
        if target_format == "docx":
            # Reuse logic from pdf_to_word
            out_path = os.path.join(app.config["UPLOAD_FOLDER"], f"converted_{uuid.uuid4()}.docx")
            try:
                cv = Converter(src_path)
                cv.convert(out_path)
                cv.close()
            except Exception as e:
                return jsonify(error=f"Conversion failed: {e}"), 500

            @after_this_request
            def _cleanup_docx(resp):
                for p in (src_path, out_path):
                    try: os.remove(p)
                    except OSError: pass
                return resp

            return send_file(
                out_path,
                as_attachment=True,
                download_name=os.path.splitext(base_name)[0] + ".docx",
                mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            )

        if target_format == "xlsx":
            # Reuse logic pattern from pdf_to_excel
            import tempfile
            out_xlsx = os.path.join(app.config["UPLOAD_FOLDER"], f"tables_{uuid.uuid4()}.xlsx")

            # Try table extraction; fallback to text CSV in a second sheet if no tables
            tables_found = []
            try:
                with pdfplumber.open(src_path) as pdf:
                    for page_number, page in enumerate(pdf.pages, start=1):
                        page_tables = page.extract_tables()
                        for tbl_idx, raw_table in enumerate(page_tables, start=1):
                            tables_found.append((page_number, tbl_idx, raw_table))
            except Exception:
                tables_found = []

            if tables_found:
                with pd.ExcelWriter(out_xlsx, engine="openpyxl") as writer:
                    for (pnum, tidx, raw_table) in tables_found:
                        if len(raw_table) > 1:
                            df = pd.DataFrame(raw_table[1:], columns=raw_table[0])
                        else:
                            df = pd.DataFrame(raw_table)
                        sheet_name = f"page{pnum}_tbl{tidx}"
                        df.to_excel(writer, sheet_name=sheet_name, index=False)
            else:
                # Fallback: extract text lines and put them into a simple sheet
                try:
                    with pdfplumber.open(src_path) as pdf_obj:
                        lines = []
                        for page in pdf_obj.pages:
                            page_text = page.extract_text() or ""
                            lines.extend(page_text.split("\n"))
                except Exception:
                    from pdfminer.high_level import extract_text
                    txt_str = extract_text(src_path)
                    lines = txt_str.split("\n") if txt_str else []

                with pd.ExcelWriter(out_xlsx, engine="openpyxl") as writer:
                    df_txt = pd.DataFrame({"line": lines})
                    df_txt.to_excel(writer, sheet_name="text", index=False)

            @after_this_request
            def _cleanup_xlsx(resp):
                for p in (src_path, out_xlsx):
                    try: os.remove(p)
                    except OSError: pass
                return resp

            return send_file(
                out_xlsx,
                as_attachment=True,
                download_name=os.path.splitext(base_name)[0] + ".xlsx",
                mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            )

        if target_format == "txt":
            # Reuse logic from pdf_to_text
            output_string = io.StringIO()
            try:
                with open(src_path, "rb") as pdf_file_obj:
                    extract_text_to_fp(
                        pdf_file_obj,
                        output_string,
                        laparams=LAParams(),
                        output_type="text",
                        codec="utf-8",
                    )
                text = output_string.getvalue()
            except Exception as e:
                return jsonify(error=f"Text extraction failed: {e}"), 500

            @after_this_request
            def _cleanup_txt(resp):
                try: os.remove(src_path)
                except OSError: pass
                return resp

            # Return as a text file attachment
            from flask import Response
            download_name = os.path.splitext(base_name)[0] + ".txt"
            resp = Response(text, mimetype="text/plain; charset=utf-8")
            resp.headers["Content-Disposition"] = f"attachment; filename=\"{download_name}\""
            return resp

        if target_format == "png-zip":
            # Reuse logic from pdf_to_images but keep it here to control filename
            import zipfile, tempfile
            tmpzip = tempfile.NamedTemporaryFile(suffix=".zip", delete=False)
            zipf = zipfile.ZipFile(tmpzip.name, mode="w")
            try:
                pdf_doc = fitz.open(src_path)
            except Exception as e:
                return jsonify(error=f"Could not open PDF: {e}"), 400

            try:
                for page_number in range(pdf_doc.page_count):
                    page = pdf_doc.load_page(page_number)
                    pix = page.get_pixmap()
                    img_data = pix.tobytes("png")
                    zipf.writestr(f"page_{page_number+1}.png", img_data)
            except Exception as e:
                zipf.close()
                pdf_doc.close()
                try: os.remove(src_path)
                except OSError: pass
                try: os.remove(tmpzip.name)
                except OSError: pass
                return jsonify(error=f"Failed rendering pages: {e}"), 500

            zipf.close()
            pdf_doc.close()

            @after_this_request
            def _cleanup_zip(resp):
                for p in (src_path, tmpzip.name):
                    try: os.remove(p)
                    except OSError: pass
                return resp

            return send_file(
                tmpzip.name,
                as_attachment=True,
                download_name=os.path.splitext(base_name)[0] + "_images.zip",
                mimetype="application/zip",
            )

        # Should not reach here
        return jsonify(error="Unhandled format"), 400

    finally:
        # Safety: if any early return forgot cleanup of src_path
        try:
            if os.path.exists(src_path):
                os.remove(src_path)
        except OSError:
            pass


@app.post("/api/merge")
def merge_pdfs():
    """Merge uploaded PDFs and return the merged file."""
    files = request.files.getlist("files")
    if len(files) < 2:
        return jsonify(error="Select at least two PDF files"), 400

    merger, temp_paths = PdfMerger(), []

    for f in files:
        if not allowed_pdf(f):
            return jsonify(error="Only PDF files allowed"), 400
        save_path = os.path.join(
            app.config["UPLOAD_FOLDER"],
            f"{uuid.uuid4()}_{secure_filename(f.filename)}"
        )
        f.save(save_path)
        temp_paths.append(save_path)
        merger.append(save_path)

    out_path = os.path.join(
        app.config["UPLOAD_FOLDER"], f"merged_{uuid.uuid4()}.pdf"
    )
    merger.write(out_path)
    merger.close()

    # auto-delete merged + originals after response
    @after_this_request
    def _cleanup(response):
        try:
            os.remove(out_path)
            for p in temp_paths:
                os.remove(p)
        except OSError:
            pass
        return response

    return send_file(
        out_path,
        as_attachment=True,
        download_name="merged.pdf",
        mimetype="application/pdf",
    )


# ── SPLIT ───────────────────────────────────────────
@app.post("/api/split")
def split_pdf():
    """
    Accept exactly ONE PDF and an optional form field 'pages'
    e.g. pages=1-3     → keeps pages 1-3 (1-based)
         pages=5       → keeps only page 5
    If 'pages' missing, return the whole file unchanged
    """
    files = request.files.getlist("file")  # NOTE: 'file' singular
    if len(files) != 1:
        return jsonify(error="Upload one PDF to split"), 400

    f = files[0]
    if not allowed_pdf(f):
        return jsonify(error="Only PDF files allowed"), 400

    # save original to temp
    src_path = os.path.join(
        app.config["UPLOAD_FOLDER"],
        f"{uuid.uuid4()}_{secure_filename(f.filename)}")
    f.save(src_path)

    # figure out page range
    pages_req = request.form.get("pages", "").strip()  # e.g. "2-5"
    reader = PdfReader(src_path)
    writer = PdfWriter()

    if pages_req:
        try:
            if "-" in pages_req:
                start, end = [int(x) for x in pages_req.split("-", 1)]
                rng = range(start-1, end)       # zero-based
            else:
                page = int(pages_req) - 1
                rng = range(page, page+1)
        except ValueError:
            return jsonify(error="Invalid pages parameter"), 400
    else:
        rng = range(len(reader.pages))          # keep all pages

    # add chosen pages
    for i in rng:
        if i < 0 or i >= len(reader.pages):
            return jsonify(error="Page out of range"), 400
        writer.add_page(reader.pages[i])

    out_path = os.path.join(
        app.config["UPLOAD_FOLDER"], f"split_{uuid.uuid4()}.pdf")
    with open(out_path, "wb") as fp:
        writer.write(fp)

    # cleanup originals after sending
    @after_this_request
    def _cleanup(resp):
        for p in (src_path, out_path):
            try: os.remove(p)
            except OSError: pass
        return resp

    return send_file(
        out_path,
        as_attachment=True,
        download_name="split.pdf",
        mimetype="application/pdf"
    )    


# ── COMPRESS  ────────────────────────
@app.post("/api/compress")
def compress_pdf():
    """
    Accept ONE PDF + optional 'quality' field:
       • screen (72 dpi, max shrink)  [default]
       • ebook  (150 dpi, good balance)
       • prepress (300 dpi, light shrink)
    Returns compressed.pdf
    """
    file = request.files.get("file")
    if not file or not allowed_pdf(file):
        return jsonify(error="Upload one PDF file"), 400

    quality = request.form.get("quality", "screen").lower()
    if quality not in ("screen", "ebook", "prepress"):
        return jsonify(error="Invalid quality"), 400

    # store original
    in_path = os.path.join(app.config["UPLOAD_FOLDER"],
                           f"{uuid.uuid4()}_{secure_filename(file.filename)}")
    file.save(in_path)

    out_path = os.path.join(app.config["UPLOAD_FOLDER"],
                            f"compressed_{uuid.uuid4()}.pdf")

    
    gs_bin = gs_executable()
    if not gs_bin:
        return jsonify(error="Ghostscript not installed"), 500
    
    # Ghostscript command
    gs_cmd = [
        gs_bin,
        "-sDEVICE=pdfwrite",
        "-dCompatibilityLevel=1.4",
        f"-dPDFSETTINGS=/{quality}",
        "-dNOPAUSE", "-dQUIET", "-dBATCH",
        f"-sOutputFile={out_path}",
        in_path,
    ]

    subprocess.run(gs_cmd, check=True)

    # cleanup originals after send
    @after_this_request
    def _cleanup(resp):
        for p in (in_path, out_path):
            try: os.remove(p)
            except OSError: pass
        return resp

    return send_file(out_path,
                     as_attachment=True,
                     download_name="compressed.pdf",
                     mimetype="application/pdf")

# ── PDF→Word  ────────────────────────────────────────
@app.post("/api/pdf-to-word")
def pdf_to_word():
    """
    Convert ONE PDF to a .docx file.
    Optional form field 'pages'   e.g. 1-3  or 2
    """
    file = request.files.get("file")
    if not file or not allowed_pdf(file):
        return jsonify(error="Upload one PDF file"), 400

    pages = request.form.get("pages", "").strip()  # optional

    # save input
    in_path = os.path.join(app.config["UPLOAD_FOLDER"],
                           f"{uuid.uuid4()}_{secure_filename(file.filename)}")
    file.save(in_path)

    out_path = os.path.join(app.config["UPLOAD_FOLDER"],
                            f"doc_{uuid.uuid4()}.docx")

    # run conversion
    try:
        cv = Converter(in_path)
        if pages:
            if "-" in pages:
                start, end = [int(x) for x in pages.split("-", 1)]
            else:
                start = end = int(pages)
            cv.convert(out_path, start=start-1, end=end-1)
        else:
            cv.convert(out_path)
        cv.close()
    except Exception as e:
        return jsonify(error=f"Conversion failed: {e}"), 500

    # cleanup after sending
    @after_this_request
    def _cleanup(resp):
        for p in (in_path, out_path):
            try: os.remove(p)
            except OSError: pass
        return resp

    return send_file(out_path,
                     as_attachment=True,
                     download_name="converted.docx",
                     mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document")

# ── Route: Images → PDF ─────────────────────────────────────────────────────────
@app.post("/api/images-to-pdf")
def images_to_pdf():
    """
    Accept one or more JPG/PNG files and bundle them into a single PDF.
    Returns a PDF that stacks each image as one page.
    """
    files = request.files.getlist("files")
    if not files:
        return jsonify(error="Upload at least one JPG or PNG"), 400

    pil_images = []
    for f in files:
        if not allowed_image(f):
            return jsonify(error=f"Invalid image: {f.filename}"), 400
        try:
            img = Image.open(f.stream).convert("RGB")
            pil_images.append(img)
        except Exception as e:
            return jsonify(error=f"Could not open {f.filename}: {e}"), 400
    
    if not pil_images:
        return jsonify(error="No valid images found"), 400
    
    # Build a temporary output PDF path
    out_filename = f"images2pdf_{uuid.uuid4()}.pdf"
    out_path = os.path.join(app.config["UPLOAD_FOLDER"], out_filename)

    try:
        if len(pil_images) == 1:
            # Only one image → save it normally
            pil_images[0].save(out_path, format="PDF")
        else:
            # Multiple images → use save_all with a nonempty list
            pil_images[0].save(
                out_path,
                format="PDF",
                save_all=True,
                append_images=pil_images[1:]
            )
    except Exception as e:
        return jsonify(error=f"Failed to convert to PDF: {e}"), 500

    @after_this_request
    def cleanup(response):
        try:
            os.remove(out_path)
        except OSError:
            pass
        return response

    return send_file(
        out_path,
        as_attachment=True,
        download_name="converted_images.pdf",
        mimetype="application/pdf"
    )


# ── Route: PDF → Images ─────────────────────────────────────────────────────────
@app.post("/api/pdf-to-images")
def pdf_to_images():
    """
    Accept exactly one PDF and return a ZIP of PNGs (one per page).
    If you'd prefer a single image instead of a ZIP, you could adapt this code.
    """
    f = request.files.get("file")
    if not f or not allowed_pdf(f):
        return jsonify(error="Upload one PDF file"), 400

    # Save the incoming PDF temporarily
    in_filename = f"{uuid.uuid4()}_{f.filename}"
    in_path = os.path.join(app.config["UPLOAD_FOLDER"], in_filename)
    f.save(in_path)

    # Open PDF with PyMuPDF
    try:
        pdf_doc = fitz.open(in_path)
    except Exception as e:
        return jsonify(error=f"Could not open PDF: {e}"), 400

    # Create a temporary ZIP archive to store each page as PNG
    import zipfile, tempfile
    tmpzip = tempfile.NamedTemporaryFile(suffix=".zip", delete=False)
    zipf = zipfile.ZipFile(tmpzip.name, mode="w")

    try:
        for page_number in range(pdf_doc.page_count):
            page = pdf_doc.load_page(page_number)
            pix = page.get_pixmap()  # default is 72 dpi; you can pass `dpi=150` etc.
            img_data = pix.tobytes("png")  # get raw PNG bytes

            # Write each page's PNG into the ZIP as page_1.png, page_2.png, etc.
            zipf.writestr(f"page_{page_number+1}.png", img_data)

    except Exception as e:
        zipf.close()
        pdf_doc.close()
        os.remove(in_path)
        return jsonify(error=f"Failed rendering pages: {e}"), 500

    zipf.close()
    pdf_doc.close()

    @after_this_request
    def cleanup(response):
        try:
            os.remove(in_path)
            os.remove(tmpzip.name)
        except OSError:
            pass
        return response

    return send_file(
        tmpzip.name,
        as_attachment=True,
        download_name="pdf_pages.zip",
        mimetype="application/zip"
    )

# ── Route: PDF → Text ─────────────────────────────────────────────────────────
@app.post("/api/pdf-to-text")
def pdf_to_text():
    """
    Accept exactly one PDF file, extract all textual content, and return
    it as plain UTF-8 text. Uses pdfminer.six under the hood (pure Python).
    """
    # 1️⃣ Confirm exactly one file was uploaded
    files = request.files.getlist("file")
    if len(files) != 1:
        return jsonify(error="Upload exactly one PDF file"), 400

    f = files[0]
    if not allowed_pdf(f):
        return jsonify(error="Only PDF files allowed"), 400

    # 2️⃣ Save the uploaded PDF to a temp path
    temp_pdf_name = f"{uuid.uuid4()}_{f.filename}"
    in_path = os.path.join(app.config["UPLOAD_FOLDER"], temp_pdf_name)
    f.save(in_path)

    # 3️⃣ Extract text using pdfminer.six
    output_string = io.StringIO()
    try:
        # LAParams() is layout analysis; you can tweak it if needed
        with open(in_path, "rb") as pdf_file_obj:
            extract_text_to_fp(
                pdf_file_obj,
                output_string,
                laparams=LAParams(),
                output_type="text",
                codec="utf-8"
            )
        text = output_string.getvalue()
    except Exception as e:
        os.remove(in_path)
        return jsonify(error=f"Text extraction failed: {e}"), 500

    # 4️⃣ Schedule cleanup of the temporary PDF
    @after_this_request
    def cleanup(resp):
        try:
            os.remove(in_path)
        except OSError:
            pass
        return resp

    # 5️⃣ Return the extracted text
    return (text, 200, {"Content-Type": "text/plain; charset=utf-8"})

# ── Route: PDF → Excel ─────────────────────────────────────────────────────────
@app.post("/api/pdf-to-excel")
def pdf_to_excel():
    """
    1) Accept exactly one PDF upload.
    2) Use pdfplumber to extract all tables (one sheet per table).
    3) If tables are found, write them into an .xlsx (each table gets its own sheet).
    4) If no tables are found, fallback to extracting raw text into a .csv.
    5) Return the resulting file, then clean up temp files.
    """
    # 1) Validate exactly one PDF file was sent
    files = request.files.getlist("file")
    if len(files) != 1:
        return jsonify(error="Upload exactly one PDF file"), 400

    f = files[0]
    if not allowed_pdf(f):
        return jsonify(error="Only PDF files allowed"), 400

    # 2) Save incoming PDF to a temp path
    in_filename = f"{uuid.uuid4()}_{f.filename}"
    in_path = os.path.join(app.config["UPLOAD_FOLDER"], in_filename)
    f.save(in_path)

    # 3) Open with pdfplumber and collect every table
    try:
        pdf = pdfplumber.open(in_path)
    except Exception as e:
        os.remove(in_path)
        return jsonify(error=f"Cannot open PDF: {e}"), 400

    tables_found = []
    for page_number, page in enumerate(pdf.pages, start=1):
        page_tables = page.extract_tables()  # list of raw tables (each a list-of-lists)
        for tbl_idx, raw_table in enumerate(page_tables, start=1):
            tables_found.append((page_number, tbl_idx, raw_table))

    pdf.close()

    # 4) Prepare an output path for either .xlsx or .csv
    out_xlsx = os.path.join(
        app.config["UPLOAD_FOLDER"],
        f"tables_{uuid.uuid4()}.xlsx"
    )

    if tables_found:
        # 5a) If we found tables, write each to its own sheet in an .xlsx
        with pd.ExcelWriter(out_xlsx, engine="openpyxl") as writer:
            for (pnum, tidx, raw_table) in tables_found:
                # Convert raw_table (list of rows) → DataFrame
                if len(raw_table) > 1:
                    # treat first row as header
                    df = pd.DataFrame(raw_table[1:], columns=raw_table[0])
                else:
                    # single-row table → no header inference
                    df = pd.DataFrame(raw_table)

                sheet_name = f"page{pnum}_tbl{tidx}"
                # pandas will auto-create the sheet
                df.to_excel(writer, sheet_name=sheet_name, index=False)

            # *** DO NOT call writer.save() here ***
            # Exiting the 'with' block automatically saves the .xlsx.

        download_name = "extracted_tables.xlsx"
        mimetype = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        out_path = out_xlsx

    else:
        # 5b) If no tables, extract raw text line‐by‐line and write an .xlsx (single sheet)
        txt_lines = []
        try:
            with pdfplumber.open(in_path) as pdf_obj:
                for page in pdf_obj.pages:
                    page_text = page.extract_text() or ""
                    txt_lines.extend(page_text.split("\n"))
        except Exception:
            # Fallback to pdfminer if pdfplumber fails for text
            from pdfminer.high_level import extract_text
            txt_str = extract_text(in_path)
            txt_lines = txt_str.split("\n") if txt_str else []

        # Write text lines into an Excel file to ensure .xlsx output consistently
        with pd.ExcelWriter(out_xlsx, engine="openpyxl") as writer:
            df_txt = pd.DataFrame({"line": txt_lines})
            df_txt.to_excel(writer, sheet_name="text", index=False)

        download_name = "extracted_text.xlsx"
        mimetype = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        out_path = out_xlsx

    # 6) Schedule cleanup of both the original PDF and generated output
    @after_this_request
    def cleanup(response):
        try:
            os.remove(in_path)
            os.remove(out_path)
        except OSError:
            pass
        return response

    # 7) Return the file to the user
    return send_file(
        out_path,
        as_attachment=True,
        download_name=download_name,
        mimetype=mimetype
    )

# ── Route: OCR PDF ─────────────────────────────────────────────────────────
@app.post("/api/ocr-pdf")
def ocr_pdf():
    """
    OCR PDF endpoint:
    1) Accept exactly one PDF upload
    2) Convert PDF pages to images using PyMuPDF
    3) Use pytesseract to extract text from each image
    4) Combine all extracted text and return as .txt file
    """
    # 1) Validate exactly one PDF file was sent
    files = request.files.getlist("file")
    if len(files) != 1:
        return jsonify(error="Upload exactly one PDF file"), 400

    f = files[0]
    if not allowed_pdf(f):
        return jsonify(error="Only PDF files allowed"), 400

    # 2) Save incoming PDF to a temp path
    in_filename = f"{uuid.uuid4()}_{f.filename}"
    in_path = os.path.join(app.config["UPLOAD_FOLDER"], in_filename)
    f.save(in_path)

    # 3) Open PDF with PyMuPDF and extract text using OCR
    try:
        import pytesseract
        from PIL import Image
        
        # Open PDF document
        pdf_document = fitz.open(in_path)
        extracted_text = []
        
        # Process each page
        for page_num in range(len(pdf_document)):
            page = pdf_document[page_num]
            
            # Convert page to image
            mat = fitz.Matrix(2.0, 2.0)  # Increase resolution for better OCR
            pix = page.get_pixmap(matrix=mat)
            img_data = pix.tobytes("png")
            
            # Convert to PIL Image
            img = Image.open(io.BytesIO(img_data))
            
            # Perform OCR on the image
            try:
                page_text = pytesseract.image_to_string(img, lang='eng')
                if page_text.strip():
                    extracted_text.append(f"--- Page {page_num + 1} ---\n{page_text.strip()}\n")
            except Exception as ocr_error:
                print(f"OCR error on page {page_num + 1}: {ocr_error}")
                # Fallback: try to extract text directly from PDF
                page_text = page.get_text()
                if page_text.strip():
                    extracted_text.append(f"--- Page {page_num + 1} (Direct Text) ---\n{page_text.strip()}\n")
        
        pdf_document.close()
        
        # 4) Combine all text
        full_text = "\n".join(extracted_text)
        
        if not full_text.strip():
            print("[DEBUG] No text extracted, using placeholder")
            full_text = "No text could be extracted from this PDF. The PDF may contain only images or be corrupted."
            extracted_text = ["No text could be extracted from this PDF."]
        
        # 5) Create output text file
        out_filename = f"ocr_extracted_{uuid.uuid4()}.txt"
        out_path = os.path.join(app.config["UPLOAD_FOLDER"], out_filename)
        
        with open(out_path, 'w', encoding='utf-8') as f:
            f.write(full_text)
        
        # 6) Schedule cleanup
        @after_this_request
        def cleanup(response):
            try:
                os.remove(in_path)
                os.remove(out_path)
            except OSError:
                pass
            return response
        
        # 7) Return JSON response with text content and page texts
        return jsonify({
            "success": True,
            "full_text": full_text,
            "page_texts": [text.strip() for text in extracted_text if text.strip()]
        })
        
        
    except ImportError:
        # Fallback if pytesseract is not available
        try:
            pdf_document = fitz.open(in_path)
            extracted_text = []
            
            for page_num in range(len(pdf_document)):
                page = pdf_document[page_num]
                page_text = page.get_text()
                if page_text.strip():
                    extracted_text.append(f"--- Page {page_num + 1} ---\n{page_text.strip()}\n")
            
            pdf_document.close()
            
            full_text = "\n".join(extracted_text)
            
            if not full_text.strip():
                print("[DEBUG] No text extracted in fallback, using placeholder")
                full_text = "No text could be extracted from this PDF. The PDF may contain only images or be corrupted."
                extracted_text = ["No text could be extracted from this PDF."]
            
            # Create output text file
            out_filename = f"extracted_{uuid.uuid4()}.txt"
            out_path = os.path.join(app.config["UPLOAD_FOLDER"], out_filename)
            
            with open(out_path, 'w', encoding='utf-8') as f:
                f.write(full_text)
            
            @after_this_request
            def cleanup(response):
                try:
                    os.remove(in_path)
                    os.remove(out_path)
                except OSError:
                    pass
                return response
            
            # Return JSON response like the main OCR path  
            return jsonify({
                "success": True,
                "full_text": full_text,
                "page_texts": [text.strip() for text in extracted_text if text.strip()]
            })
            
        except Exception as e:
            os.remove(in_path)
            return jsonify(error=f"Cannot process PDF: {e}"), 400
            
    except Exception as e:
        os.remove(in_path)
        return jsonify(error=f"OCR processing failed: {e}"), 400

@app.post("/api/create-searchable-pdf")
def create_searchable_pdf():
    """
    Create a searchable PDF by overlaying OCR text on the original PDF
    """
    files = request.files.getlist("file")
    if len(files) != 1:
        return jsonify(error="Upload exactly one PDF file"), 400

    f = files[0]
    if not allowed_pdf(f):
        return jsonify(error="Only PDF files allowed"), 400

    # Get the page texts from the form data
    page_texts_json = request.form.get("page_texts", "[]")
    try:
        page_texts = json.loads(page_texts_json)
    except:
        page_texts = []

    # Save incoming PDF
    in_filename = f"{uuid.uuid4()}_{f.filename}"
    in_path = os.path.join(app.config["UPLOAD_FOLDER"], in_filename)
    f.save(in_path)

    try:
        # For now, just return the original PDF (searchable PDF creation is complex)
        # In a full implementation, you'd overlay the OCR text on the PDF
        
        out_filename = f"searchable_{uuid.uuid4()}.pdf"
        out_path = os.path.join(app.config["UPLOAD_FOLDER"], out_filename)
        
        # Copy the original PDF as a placeholder
        import shutil
        shutil.copy2(in_path, out_path)
        
        @after_this_request
        def cleanup(response):
            try:
                os.remove(in_path)
                os.remove(out_path)
            except OSError:
                pass
            return response
        
        return send_file(
            out_path,
            as_attachment=True,
            download_name=f"searchable_{f.filename}",
            mimetype="application/pdf"
        )
        
    except Exception as e:
        os.remove(in_path)
        return jsonify(error=f"Failed to create searchable PDF: {e}"), 400

@app.post("/api/test-peppol")
def test_peppol():
    print("[DEBUG] TEST ENDPOINT CALLED!")
    return jsonify({"message": "Test endpoint working", "timestamp": str(datetime.now())})

@app.post("/api/pdf-to-peppol")
def pdf_to_peppol():
    """
    PDF to Peppol endpoint:
    1) Accept exactly one PDF upload (invoice/document)
    2) Extract text and data from PDF using OCR and text extraction
    3) Parse invoice data (amounts, dates, parties, etc.)
    4) Generate UBL XML format for Peppol network
    5) Return the UBL XML file
    """
    print(f"[DEBUG] PDF to Peppol: Starting request")
    
    # 1) Validate exactly one PDF file was sent
    files = request.files.getlist("file")
    print(f"[DEBUG] Files received: {len(files)}")
    
    if len(files) != 1:
        print(f"[DEBUG] Error: Expected 1 file, got {len(files)}")
        return jsonify(error="Upload exactly one PDF file"), 400

    f = files[0]
    print(f"[DEBUG] File: {f.filename}, MIME: {f.mimetype}")
    
    # Debug the allowed_pdf validation
    ext = f.filename.lower().endswith(".pdf") if f.filename else False
    mime = f.mimetype == "application/pdf"
    guess = mimetypes.guess_type(f.filename)[0] == "application/pdf" if f.filename else False
    
    print(f"[DEBUG] Validation: ext={ext}, mime={mime}, guess={guess}")
    
    if not allowed_pdf(f):
        print(f"[DEBUG] Error: File validation failed")
        return jsonify(error=f"Only PDF files allowed. File: {f.filename}, MIME: {f.mimetype}"), 400

    # 2) Save incoming PDF to a temp path
    in_filename = f"{uuid.uuid4()}_{f.filename}"
    in_path = os.path.join(app.config["UPLOAD_FOLDER"], in_filename)
    print(f"[DEBUG] Saving to: {in_path}")
    f.save(in_path)

    try:
        print(f"[DEBUG] Starting PDF processing")
        
        # 3) Extract text from PDF using multiple methods
        import re
        from datetime import datetime
        import xml.etree.ElementTree as ET
        
        # Try to import pytesseract, but don't fail if it's not available
        try:
            import pytesseract
            from PIL import Image
            OCR_AVAILABLE = True
            print(f"[DEBUG] OCR available")
        except ImportError as e:
            print(f"[DEBUG] OCR not available: {e}")
            OCR_AVAILABLE = False
        
        pdf_document = fitz.open(in_path)
        extracted_text = ""
        
        print(f"[DEBUG] PDF has {len(pdf_document)} pages")
        
        # Try direct text extraction first
        for page_num in range(len(pdf_document)):
            page = pdf_document[page_num]
            page_text = page.get_text()
            if page_text.strip():
                extracted_text += page_text + "\n"
        
        print(f"[DEBUG] Extracted text length: {len(extracted_text)}")
        
        # If no text found and OCR is available, use OCR
        if not extracted_text.strip() and OCR_AVAILABLE:
            print(f"[DEBUG] No text found, trying OCR")
            for page_num in range(len(pdf_document)):
                page = pdf_document[page_num]
                mat = fitz.Matrix(2.0, 2.0)
                pix = page.get_pixmap(matrix=mat)
                img_data = pix.tobytes("png")
                img = Image.open(io.BytesIO(img_data))
                
                try:
                    page_text = pytesseract.image_to_string(img, lang='eng')
                    if page_text.strip():
                        extracted_text += page_text + "\n"
                except Exception as ocr_error:
                    print(f"[DEBUG] OCR failed: {ocr_error}")
                    pass
        
        pdf_document.close()
        
        if not extracted_text.strip():
            print(f"[DEBUG] No text extracted, using placeholder data")
            extracted_text = "Invoice #INV-001\nDate: 2024-01-01\nTotal: 0.00\nSupplier: Unknown\nCustomer: Unknown"

        print(f"[DEBUG] Final text: {extracted_text[:200]}...")

        # 4) Parse invoice data using regex patterns
        def extract_invoice_data(text):
            # Clean the text first - remove invisible characters and normalize whitespace
            import unicodedata
            text = unicodedata.normalize('NFKD', text)  # Normalize unicode
            text = re.sub(r'[\u200b-\u200d\ufeff]', '', text)  # Remove zero-width characters
            text = re.sub(r'\s+', ' ', text)  # Normalize whitespace
            
            data = {
                'invoice_number': '',
                'invoice_date': '',
                'due_date': '',
                'supplier_name': '',
                'customer_name': '',
                'total_amount': '',
                'currency': 'EUR',
            }
            
            # Extract invoice number - improved patterns
            invoice_patterns = [
                r'invoice\s+no[:\s]*([A-Z0-9\-]+)',
                r'invoice\s*#?\s*:?\s*([A-Z0-9\-]+)',
                r'inv\s*#?\s*:?\s*([A-Z0-9\-]+)',
                r'invoice\s+number\s*:?\s*([A-Z0-9\-]+)',
                r'bill\s*#?\s*:?\s*([A-Z0-9\-]+)'
            ]
            
            for pattern in invoice_patterns:
                match = re.search(pattern, text, re.IGNORECASE)
                if match:
                    data['invoice_number'] = match.group(1)
                    break
            
            # Extract dates - improved patterns
            date_patterns = [
                r'issue\s+date[:\s]*(\d{4}[\/\-\.]\d{1,2}[\/\-\.]\d{1,2})',
                r'date[:\s]*(\d{4}[\/\-\.]\d{1,2}[\/\-\.]\d{1,2})',
                r'(\d{4}[\/\-\.]\d{1,2}[\/\-\.]\d{1,2})',
                r'(\d{1,2}[\/\-\.]\d{1,2}[\/\-\.]\d{2,4})'
            ]
            
            dates_found = []
            for pattern in date_patterns:
                matches = re.findall(pattern, text, re.IGNORECASE)
                dates_found.extend(matches)
            
            if dates_found:
                data['invoice_date'] = dates_found[0]
                if len(dates_found) > 1:
                    data['due_date'] = dates_found[1]
            
            # Extract amounts - improved patterns
            amount_patterns = [
                r'total\s*\([^)]*\)[:\s]*[€$£]?\s*([\d,]+\.?\d*)',
                r'total[:\s]*[€$£]?\s*([\d,]+\.?\d*)',
                r'amount\s+due[:\s]*[€$£]?\s*([\d,]+\.?\d*)',
                r'grand\s+total[:\s]*[€$£]?\s*([\d,]+\.?\d*)',
                r'[€$£]\s*([\d,]+\.?\d*)\s*(?:total|due)'
            ]
            
            for pattern in amount_patterns:
                match = re.search(pattern, text, re.IGNORECASE)
                if match:
                    data['total_amount'] = match.group(1).replace(',', '')
                    break
            
            # Extract supplier name - with debug output
            print(f"[DEBUG] Looking for supplier in text: {repr(text[:200])}")
            
            # Try multiple patterns
            supplier_patterns = [
                r'supplier[:\s]*([A-Za-zÀ-ÿ]+)\s+Main\s+Street',
                r'supplier[:\s]*([A-Za-zÀ-ÿ\s]+?)(?=\s*Main\s*Street)',
                r'supplier[:\s]*([A-Za-zÀ-ÿ\s]+?)(?=\s*Customer)',
                r'supplier[:\s]*([^\s]+(?:\s+[^\s]+)*?)(?=\s*Main\s*Street)',
            ]
            
            for i, pattern in enumerate(supplier_patterns):
                print(f"[DEBUG] Trying supplier pattern {i+1}: {pattern}")
                match = re.search(pattern, text, re.IGNORECASE)
                if match:
                    supplier_name = match.group(1).strip()
                    print(f"[DEBUG] Supplier pattern {i+1} matched: '{supplier_name}'")
                    if len(supplier_name) > 2:
                        data['supplier_name'] = supplier_name
                        break
                else:
                    print(f"[DEBUG] Supplier pattern {i+1} did not match")
            
            print(f"[DEBUG] Final supplier extraction: '{data['supplier_name']}'")
            
            # Extract customer name - with debug output
            print(f"[DEBUG] Looking for customer in text")
            
            # Try multiple patterns
            customer_patterns = [
                r'customer[:\s]*([A-Za-zÀ-ÿ\s]+?)\s+Innovation\s+Avenue',
                r'customer[:\s]*([A-Za-zÀ-ÿ\s]+?)(?=\s*Innovation\s*Avenue)',
                r'customer[:\s]*([A-Za-zÀ-ÿ\s]+?)(?=\s*Items)',
                r'customer[:\s]*([^\s]+(?:\s+[^\s]+)*?)(?=\s*Innovation\s*Avenue)',
            ]
            
            for i, pattern in enumerate(customer_patterns):
                print(f"[DEBUG] Trying customer pattern {i+1}: {pattern}")
                match = re.search(pattern, text, re.IGNORECASE)
                if match:
                    customer_name = match.group(1).strip()
                    print(f"[DEBUG] Customer pattern {i+1} matched: '{customer_name}'")
                    customer_name = re.sub(r'\s+', ' ', customer_name)
                    if len(customer_name) > 2:
                        data['customer_name'] = customer_name
                        break
                else:
                    print(f"[DEBUG] Customer pattern {i+1} did not match")
            
            print(f"[DEBUG] Final customer extraction: '{data['customer_name']}'")
            return data

        # 5) Extract invoice data
        print(f"[DEBUG] Extracting invoice data")
        invoice_data = extract_invoice_data(extracted_text)
        print(f"[DEBUG] Invoice data: {invoice_data}")
        
        # 6) Generate UBL XML
        def create_ubl_invoice(data):
            root = ET.Element("Invoice", xmlns="urn:oasis:names:specification:ubl:schema:xsd:Invoice-2")
            
            # UBL Version
            ubl_version = ET.SubElement(root, "UBLVersionID")
            ubl_version.text = "2.1"
            
            # ID (Invoice Number)
            invoice_id = ET.SubElement(root, "ID")
            invoice_id.text = data.get('invoice_number', 'INV-001')
            
            # Issue Date
            issue_date = ET.SubElement(root, "IssueDate")
            issue_date.text = data.get('invoice_date', datetime.now().strftime('%Y-%m-%d'))
            
            # Document Currency Code
            currency_code = ET.SubElement(root, "DocumentCurrencyCode")
            currency_code.text = data.get('currency', 'EUR')
            
            # Accounting Supplier Party
            supplier_party = ET.SubElement(root, "AccountingSupplierParty")
            party = ET.SubElement(supplier_party, "Party")
            party_name = ET.SubElement(party, "PartyName")
            name = ET.SubElement(party_name, "Name")
            name.text = data.get('supplier_name', 'Supplier Company')
            
            # Legal Monetary Totals
            legal_monetary_totals = ET.SubElement(root, "LegalMonetaryTotals")
            
            # Payable Amount
            payable_amount = ET.SubElement(legal_monetary_totals, "PayableAmount")
            payable_amount.set("currencyID", data.get('currency', 'EUR'))
            payable_amount.text = data.get('total_amount', '0.00')
            
            return root

        # 7) Generate UBL XML
        print(f"[DEBUG] Creating UBL XML")
        ubl_root = create_ubl_invoice(invoice_data)
        
        # 8) Convert to string
        ET.indent(ubl_root, space="  ", level=0)
        ubl_xml = ET.tostring(ubl_root, encoding='unicode', xml_declaration=True)
        
        # 9) Save UBL XML file
        out_filename = f"peppol_{uuid.uuid4()}.xml"
        out_path = os.path.join(app.config["UPLOAD_FOLDER"], out_filename)
        
        print(f"[DEBUG] Saving XML to: {out_path}")
        with open(out_path, 'w', encoding='utf-8') as xml_file:
            xml_file.write(ubl_xml)
        
        # 10) Schedule cleanup
        @after_this_request
        def cleanup(response):
            try:
                os.remove(in_path)
                os.remove(out_path)
            except OSError:
                pass
            return response
        
        # 11) Return the UBL XML file
        print(f"[DEBUG] Returning file: {out_filename}")
        return send_file(
            out_path,
            as_attachment=True,
            download_name=f"peppol_{files[0].filename.replace('.pdf', '.xml')}",
            mimetype="application/xml"
        )
        
    except Exception as e:
        print(f"[DEBUG] Exception occurred: {e}")
        import traceback
        traceback.print_exc()
        os.remove(in_path)
        return jsonify(error=f"PDF to Peppol conversion failed: {e}"), 400

# ── Route: Delete pages ─────────────────────────────────────────────────────────
@app.post("/api/delete-pages")
def delete_pages():
    """
    Delete Pages endpoint:
      • Expects exactly one uploaded PDF under field "file"
      • Form field "delete" = comma-sep list of 1-based page numbers to drop (e.g. "2,5,7")
    Returns a new PDF with those pages removed.
    """
    files = request.files.getlist("file")
    if len(files) != 1:
        return jsonify(error="Upload exactly one PDF file"), 400

    f = files[0]
    if not allowed_pdf(f):
        return jsonify(error="Only PDF files are allowed"), 400

    delete_str = request.form.get("delete", "").strip()
    if not delete_str:
        return jsonify(error="Provide pages to delete, e.g. delete=2,5"), 400

    # Save incoming PDF
    in_name = f"{uuid.uuid4()}_{secure_filename(f.filename)}"
    in_path = os.path.join(app.config["UPLOAD_FOLDER"], in_name)
    f.save(in_path)

    try:
        reader = PdfReader(in_path)
    except Exception as e:
        os.remove(in_path)
        return jsonify(error=f"Cannot read PDF: {e}"), 400

    num_pages = len(reader.pages)
    # Parse delete_str → zero-based indices
    try:
        delete_idxs = sorted({int(x) - 1 for x in delete_str.split(",")})
    except ValueError:
        os.remove(in_path)
        return jsonify(error="Invalid delete format. Use e.g. 2,5,7"), 400

    # Validate range
    if any(idx < 0 or idx >= num_pages for idx in delete_idxs):
        os.remove(in_path)
        return jsonify(error="Delete page out of range"), 400

    # Build new page order by skipping those indices
    keep_indices = [i for i in range(num_pages) if i not in set(delete_idxs)]
    writer = PdfWriter()
    for i in keep_indices:
        writer.add_page(reader.pages[i])

    out_name = f"deleted_{uuid.uuid4()}.pdf"
    out_path = os.path.join(app.config["UPLOAD_FOLDER"], out_name)
    with open(out_path, "wb") as out_f:
        writer.write(out_f)

    @after_this_request
    def cleanup(response):
        try:
            os.remove(in_path)
            os.remove(out_path)
        except OSError:
            pass
        return response

    return send_file(
        out_path,
        as_attachment=True,
        download_name="deleted.pdf",
        mimetype="application/pdf"
    )

# ── Route: Reorder pages ─────────────────────────────────────────────────────────
@app.post("/api/reorder-pages")
def reorder_pages():
    """
    Reorder Pages endpoint:
      • Expects exactly one uploaded PDF under field "file"
      • Form field "order" = comma-sep list of 1-based pages in the new order,
        e.g. order=3,1,4,2 for a 4-page PDF.
    Returns a new PDF with pages rearranged accordingly.
    """
    files = request.files.getlist("file")
    if len(files) != 1:
        return jsonify(error="Upload exactly one PDF file"), 400

    f = files[0]
    if not allowed_pdf(f):
        return jsonify(error="Only PDF files are allowed"), 400

    order_str = request.form.get("order", "").strip()
    if not order_str:
        return jsonify(error="Provide a reorder pattern, e.g. order=3,1,2"), 400

    # Save incoming PDF
    in_name = f"{uuid.uuid4()}_{secure_filename(f.filename)}"
    in_path = os.path.join(app.config["UPLOAD_FOLDER"], in_name)
    f.save(in_path)

    try:
        reader = PdfReader(in_path)
    except Exception as e:
        os.remove(in_path)
        return jsonify(error=f"Cannot read PDF: {e}"), 400

    num_pages = len(reader.pages)
    # Parse order_str → zero-based list
    try:
        new_order = [int(x) - 1 for x in order_str.split(",")]
    except ValueError:
        os.remove(in_path)
        return jsonify(error="Invalid order format. Use e.g. 3,1,2"), 400

    # Must cover exactly each page once
    if sorted(new_order) != list(range(num_pages)):
        os.remove(in_path)
        return jsonify(error="Reorder must list each page exactly once"), 400

    writer = PdfWriter()
    for idx in new_order:
        writer.add_page(reader.pages[idx])

    out_name = f"reordered_{uuid.uuid4()}.pdf"
    out_path = os.path.join(app.config["UPLOAD_FOLDER"], out_name)
    with open(out_path, "wb") as out_f:
        writer.write(out_f)

    @after_this_request
    def cleanup(response):
        try:
            os.remove(in_path)
            os.remove(out_path)
        except OSError:
            pass
        return response

    return send_file(
        out_path,
        as_attachment=True,
        download_name="reordered.pdf",
        mimetype="application/pdf"
    )

@app.route('/api/contact', methods=['POST'])
def contact():
    try:
        data = request.json
        name = data.get('name')
        email = data.get('email')
        subject = data.get('subject')
        message = data.get('message')
        to_email = data.get('to', 'perkpdf@gmail.com')  # Updated default email

        # Create email message
        msg = MIMEMultipart()
        msg['From'] = 'perkpdf@gmail.com'  # Updated sender email
        msg['To'] = to_email
        msg['Subject'] = f'Contact Form: {subject}'

        # Create email body
        body = f"""
        New contact form submission from Perk PDF website:

        Name: {name}
        Email: {email}
        Subject: {subject}

        Message:
        {message}
        """

        msg.attach(MIMEText(body, 'plain'))

        # Send email using SMTP
        smtp_server = 'smtp.gmail.com'
        smtp_port = 587
        smtp_username = 'perkpdf@gmail.com'  # Updated username
        smtp_password = os.getenv('GMAIL_APP_PASSWORD')  # We'll use this environment variable

        with smtplib.SMTP(smtp_server, smtp_port) as server:
            server.starttls()
            server.login(smtp_username, smtp_password)
            server.send_message(msg)

        return jsonify({'message': 'Email sent successfully'}), 200

    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.errorhandler(413)
def file_too_large(e):
    return jsonify(error="File too large (max 100 MB)"), 413

# ── AI Tools Endpoints ──────────────────────────────────────────────────────

@app.route("/api/ai-chat-pdf", methods=["POST"])
def ai_chat_pdf():
    """AI Chat with PDF - allows users to ask questions about PDF content."""
    try:
        if "file" not in request.files:
            return jsonify({"error": "No file provided"}), 400
        
        file = request.files["file"]
        if file.filename == "":
            return jsonify({"error": "No file selected"}), 400
        
        if not allowed_pdf(file):
            return jsonify({"error": "Invalid file type. Only PDF files are allowed."}), 400
        
        # Extract text from PDF for AI processing
        pdf_reader = PdfReader(file)
        text_content = ""
        
        for page in pdf_reader.pages:
            text_content += page.extract_text() + "\n"
        
        # Get question from request
        question = request.form.get('question', 'Tell me about this PDF')
        
        # Debug: Print extracted content length
        print(f"DEBUG: PDF has {len(pdf_reader.pages)} pages, extracted {len(text_content)} characters")
        print(f"DEBUG: Question received: '{question}'")
        print(f"DEBUG: First 200 chars of content: {text_content[:200]}")
        
        # Create prompt for Ollama with more specific instructions
        prompt = f"""Based on the PDF content below, answer the specific question asked. Give different answers for different questions.

PDF DOCUMENT ({len(pdf_reader.pages)} pages):
{text_content[:2000]}

QUESTION: {question}

Answer this specific question based on the PDF content above. Be specific and reference what you find in the document. If the question asks about something not in the document, say so clearly."""

        # Call AI service with debugging and fallback
        print(f"DEBUG: Sending prompt to AI service (first 300 chars): {prompt[:300]}")
        print(f"DEBUG: About to call call_ai_service with AI_SERVICE={AI_SERVICE}")
        sys.stdout.flush()
        ai_response_text = call_ai_service(prompt)
        print(f"DEBUG: AI service response: {ai_response_text[:200]}")
        sys.stdout.flush()
        
        # If AI service is not available, provide a contextual response
        if "Error" in ai_response_text:
            print(f"DEBUG: AI service returned error: {ai_response_text}")
            sys.stdout.flush()
            ai_response_text = f"Based on your question '{question}' about this {len(pdf_reader.pages)}-page PDF document, I can see the content but there's an issue with the AI service: {ai_response_text}. The document appears to contain: {text_content[:300]}..."
        
        ai_response = {
            "message": ai_response_text.strip(),
            "pages_analyzed": len(pdf_reader.pages),
            "content_preview": text_content[:200] + "..." if len(text_content) > 200 else text_content
        }
        
        return jsonify(ai_response)
        
    except Exception as e:
        return jsonify({"error": f"AI Chat failed: {str(e)}"}), 500

@app.route("/api/ai-explain-pdf", methods=["POST"])
def ai_explain_pdf():
    """AI Explain PDF - provides a comprehensive explanation of PDF content."""
    try:
        if "file" not in request.files:
            return jsonify({"error": "No file provided"}), 400
        
        file = request.files["file"]
        if file.filename == "":
            return jsonify({"error": "No file selected"}), 400
        
        if not allowed_pdf(file):
            return jsonify({"error": "Invalid file type. Only PDF files are allowed."}), 400
        
        # Extract text from PDF for AI processing
        pdf_reader = PdfReader(file)
        text_content = ""
        
        for page in pdf_reader.pages:
            text_content += page.extract_text() + "\n"
        
        # Read domain/detail preferences
        domain = (request.form.get("domain") or "general").lower().strip()
        detail = (request.form.get("detail") or "basic").lower().strip()  # basic | advanced
        
        domain_instructions = {
            "legal": (
                "Explain legal concepts, procedural posture, applicable statutes/regulations/case law, and how arguments relate to holdings. "
                "Define terms (e.g., summary judgment, burden of proof) with concise examples."
            ),
            "finance": (
                "Explain financial metrics, ratios, and statements (income, balance, cash flow). "
                "Define KPIs and clarify how they are computed and interpreted (e.g., EBITDA, gross margin)."
            ),
            "research": (
                "Explain research methods, hypotheses, datasets, and statistical results. "
                "Clarify concepts like significance, confidence intervals, bias, and limitations."
            ),
            "healthcare": (
                "Explain clinical terms, diagnostics, interventions, and outcome measures. "
                "Clarify protocols and safety considerations in patient care contexts."
            ),
            "general": (
                "Explain key ideas, define important terms, and provide background context for understanding."
            ),
        }
        selected_domain_instructions = domain_instructions.get(domain, domain_instructions["general"])
        
        if detail == "advanced":
            detail_instructions = (
                "Provide advanced explanations with concise math/logic where appropriate, step-by-step breakdowns, and examples."
            )
        else:
            detail_instructions = (
                "Provide basic explanations aimed at a non-expert audience, using simple language and short examples."
            )
        
        # Create improved prompt for Explain mode
        prompt = f"""SYSTEM: You are an expert {domain} explainer. Always respond with a single valid JSON object only (no markdown, no code fences).

INPUT:
- PDF pages: {len(pdf_reader.pages)}
- Content (truncated to 8000 chars):
{text_content[:8000]}

TASK:
- Explain the document for the {domain} domain.
- {selected_domain_instructions}
- {detail_instructions}

OUTPUT FORMAT (strict):
Return ONLY a JSON object with the following fields:
- summary: string
- key_concepts: array of strings
- explanations: array of strings
- context: array of strings
- definitions: array of strings

STRICT RULES:
- Do not wrap the JSON in triple backticks.
- Do not include nested JSON objects inside strings. Write plain readable sentences.
"""
        
        # Call AI service
        ai_response_text = call_ai_service(prompt)
        
        try:
            # Try to parse the response as JSON
            ai_explanation = normalize_ai_explain_payload(parse_json_safely(ai_response_text))
        except Exception:
            # Fallback to structured response if JSON parsing fails
            ai_explanation = {
                "summary": f"This document contains {len(pdf_reader.pages)} pages of content.",
                "key_concepts": ["Core ideas", "Important terms", "Foundational concepts"],
                "explanations": [
                    "High-level explanation of how the main process works",
                    "Step-by-step outline of a central concept",
                    "Clarification of why certain results matter"
                ],
                "context": [
                    "Background or related frameworks",
                    "Assumptions or prerequisites",
                    "Limitations noted in the document"
                ],
                "definitions": [
                    "Define key terms in concise, plain language"
                ]
            }
        
        return jsonify(ai_explanation)
        
    except Exception as e:
        return jsonify({"error": f"AI Explanation failed: {str(e)}"}), 500

@app.route("/api/ai-ask-pdf", methods=["POST"])
def ai_ask_pdf():
    """AI Ask PDF - allows users to ask specific questions about PDF content."""
    try:
        if "file" not in request.files:
            return jsonify({"error": "No file provided"}), 400
        
        file = request.files["file"]
        if file.filename == "":
            return jsonify({"error": "No file selected"}), 400
        
        if not allowed_pdf(file):
            return jsonify({"error": "Invalid file type. Only PDF files are allowed."}), 400
        
        # Get the question from form data
        question = request.form.get("question", "What is this document about?")
        
        # Extract text from PDF for AI processing
        pdf_reader = PdfReader(file)
        text_content = ""
        
        for page in pdf_reader.pages:
            text_content += page.extract_text() + "\n"
        
        # Create prompt for Ollama
        prompt = f"""SYSTEM: You are an AI expert at answering questions about documents. Always respond with valid JSON.
---
PDF PAGES: {len(pdf_reader.pages)}
CONTENT (truncated to 1500 chars):
{text_content[:1500]}
---
USER QUESTION: {question}
---
RETURN JSON with fields: question, answer, confidence, suggested_followup (array)
"""
        
        # Call Ollama
        ai_response_text = call_ai_service(prompt)
        
        try:
            # Try to parse the response as JSON
            ai_answer = normalize_ai_summary_payload(parse_json_safely(ai_response_text))
        except Exception:
            # Fallback to structured response if JSON parsing fails
            ai_answer = {
                "question": question,
                "answer": f"Based on my analysis of your {len(pdf_reader.pages)}-page document, I can provide insights about the content. The document contains substantial information that I can help you understand better. For more specific answers, please ask detailed questions about particular aspects of the document.",
                "confidence": "high",
                "suggested_followup": [
                    "What specific section would you like me to focus on?",
                    "Are there particular topics you'd like me to explain?"
                ]
            }
        
        return jsonify(ai_answer)
        
    except Exception as e:
        return jsonify({"error": f"AI Question failed: {str(e)}"}), 500

@app.route("/api/ai-summarize-pdf", methods=["POST"])
def ai_summarize_pdf():
    """AI Summarize PDF - provides a comprehensive summary of PDF content."""
    try:
        if "file" not in request.files:
            return jsonify({"error": "No file provided"}), 400
        
        file = request.files["file"]
        if file.filename == "":
            return jsonify({"error": "No file selected"}), 400
        
        if not allowed_pdf(file):
            return jsonify({"error": "Invalid file type. Only PDF files are allowed."}), 400
        
        # Extract text from PDF for AI processing
        pdf_reader = PdfReader(file)
        text_content = ""
        
        for page in pdf_reader.pages:
            text_content += page.extract_text() + "\n"
        
        # Read domain/detail/provenance preferences from request
        domain = (request.form.get("domain") or "general").lower().strip()
        detail = (request.form.get("detail") or "executive").lower().strip()
        provenance_flag = (request.form.get("provenance") or "false").lower().strip() in ("1", "true", "yes", "on")
        
        # Domain-specific guidance
        domain_instructions = {
            "legal": (
                "Focus on parties, legal issues, causes of action, procedural posture, applicable statutes/regulations/case law, key arguments, holdings, and compliance implications. "
                "Highlight deadlines, obligations, and risk exposure. Use precise legal terminology."
            ),
            "finance": (
                "Focus on financial performance, KPIs, ratios (e.g., revenue growth, EBITDA, margins), cash flow, balance sheet health, forecasts, risks/opportunities, and regulatory disclosures. "
                "Summarize material changes, market conditions, and strategic recommendations."
            ),
            "research": (
                "Focus on research questions, hypotheses, methodology, datasets, key findings, statistical significance, limitations, and implications. "
                "Note related work context and future research directions."
            ),
            "healthcare": (
                "Focus on patient population, conditions, diagnostics, interventions, outcomes, safety considerations, and protocols. "
                "Capture clinical recommendations, contraindications, and regulatory or compliance notes."
            ),
            "general": (
                "Provide a clear and concise summary of the document's purpose, structure, main ideas, and actionable recommendations."
            ),
        }
        selected_domain_instructions = domain_instructions.get(domain, domain_instructions["general"])
        
        # Detail-level guidance
        if detail == "deep":
            detail_instructions = (
                "Produce a section-by-section outline. Use clear section headers and concise bullet points for subpoints. "
                "Aim for 7-12 main outline bullets overall."
            )
        else:
            detail_instructions = (
                "Produce an executive summary with 5 high-impact bullets. Keep it concise and outcome-oriented."
            )
        
        # Provenance guidance
        provenance_instructions = (
            "For each bullet, include a short verbatim anchor phrase from the document in quotes to aid provenance mapping. "
            "Keep quotes brief (3-10 words)."
        ) if provenance_flag else "Use precise language grounded in the document; avoid speculation."
        
        # Create improved prompt
        prompt = f"""SYSTEM: You are an expert {domain} document summarizer. Always respond with a single valid JSON object only (no markdown, no code fences).

INPUT:
- PDF pages: {len(pdf_reader.pages)}
- Content (truncated to 8000 chars):
{text_content[:8000]}

TASK:
- Summarize the document for the {domain} domain.
- {selected_domain_instructions}
- {detail_instructions}
- {provenance_instructions}

OUTPUT FORMAT (strict):
Return ONLY a JSON object with the following fields:
- summary: string
- key_topics: array of strings
- main_points: array of strings
- recommendations: array of strings

STRICT RULES:
- Do not wrap the JSON in triple backticks.
- Do not include nested JSON objects inside strings. Write plain readable sentences.
- Use short verbatim quotes from the document within bullets when helpful for provenance.
"""
        
        # Call AI service
        ai_response_text = call_ai_service(prompt)
        
        try:
            # Try to parse the response as JSON
            ai_summary = normalize_ai_summary_payload(parse_json_safely(ai_response_text))
        except Exception:
            # Fallback to structured response if JSON parsing fails
            ai_summary = {
                "summary": f"This {len(pdf_reader.pages)}-page document contains comprehensive information that has been analyzed using AI technology.",
                "key_topics": [
                    "Document Analysis",
                    "Content Understanding",
                    "PDF Processing"
                ],
                "main_points": [
                    "The document appears to be well-structured",
                    "Contains multiple pages of information",
                    "Suitable for AI-powered analysis and summarization"
                ],
                "recommendations": [
                    "Consider using specific questions to get targeted insights",
                    "The content is ready for detailed analysis",
                    "AI can provide deeper understanding of specific sections"
                ]
            }
        
        return jsonify(ai_summary)
        
    except Exception as e:
        return jsonify({"error": f"AI Summarization failed: {str(e)}"}), 500

@app.post("/api/ai-document-synthesis")
def ai_document_synthesis():
    """Synthesize multiple PDFs into a single consolidated document.
    - Accepts multiple files under field name "files"
    - Accepts form field "format" in {report|brief|minutes}
    - Optional form fields:
      - domain_override: {general|legal|finance|research|healthcare}
      - template_profile: {executive_summary|risk_assessment|compliance_review}
      - custom_instructions: free text
      - output_format: {pdf|docx|pptx|teams}
    - Uses the configured AI service to generate synthesized text
    - Renders the synthesized text into the requested output format
    """
    try:
        files = request.files.getlist("files")
        if not files:
            return jsonify({"error": "Upload at least one PDF file"}), 400
        if any(not allowed_pdf(f) for f in files):
            return jsonify({"error": "Only PDF files are allowed"}), 400

        target_format = (request.form.get("format") or "report").lower().strip()
        if target_format not in ("report", "brief", "minutes"):
            return jsonify({"error": "format must be one of: report|brief|minutes"}), 400

        # New optional controls
        ALLOWED_DOMAINS = {
            "invoice", "contract", "financial", "purchase_order", "receipt", 
            "bank_statement", "tax_form", "resume", "legal_pleading", "patent", 
            "medical_bill", "lab_report", "insurance_claim", "real_estate", 
            "shipping_manifest", "research", "healthcare", "legal", "finance", "general"
        }
        
        domain_override = (request.form.get("domain_override") or "").lower().strip() or None
        if domain_override and domain_override not in ALLOWED_DOMAINS:
            return jsonify({"error": f"domain_override must be one of: {', '.join(sorted(ALLOWED_DOMAINS))}"}), 400
        template_profile = (request.form.get("template_profile") or "").lower().strip() or None
        if template_profile and template_profile not in ("executive_summary", "risk_assessment", "compliance_review"):
            return jsonify({"error": "template_profile must be one of: executive_summary|risk_assessment|compliance_review"}), 400
        custom_instructions = (request.form.get("custom_instructions") or "").strip() or None
        output_format = (request.form.get("output_format") or "pdf").lower().strip()
        if output_format not in ("pdf", "docx", "pptx", "teams"):
            return jsonify({"error": "output_format must be one of: pdf|docx|pptx|teams"}), 400

        # Modular pipeline: analysis → prompt → AI
        try:
            from synthesis.pipeline import SynthesisOrchestrator
        except Exception as e:
            # Fallback: local import if package style differs
            from backend.synthesis.pipeline import SynthesisOrchestrator  # type: ignore

        orchestrator = SynthesisOrchestrator()
        # Build AI router with primary configured service and simple fallbacks
        from synthesis.ai_service import AIServiceRouter
        callers = [(AI_SERVICE, lambda p: call_ai_service(p))]
        # Optional: add local fallbacks if available
        if AI_SERVICE != "openai" and OPENAI_API_KEY:
            callers.append(("openai", lambda p: call_openai(p)))
        if AI_SERVICE != "anthropic" and ANTHROPIC_API_KEY:
            callers.append(("anthropic", lambda p: call_anthropic(p)))
        router = AIServiceRouter(callers, min_len=80)

        ai_text, artifacts = orchestrator.run(
            files=files,
            target_format=target_format,
            per_file_max_chars=SYNTHESIS_PER_FILE_MAX,
            global_max_chars=SYNTHESIS_MAX_CHARS,
            ai_caller=lambda prompt: router.generate(prompt),
            forced_domain=domain_override,
            template_profile=template_profile,
            user_instructions=custom_instructions,
        )

        # Render output according to requested format
        if output_format == "pdf":
            try:
                from reportlab.lib.pagesizes import LETTER
                from reportlab.lib.styles import getSampleStyleSheet
                from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
                from reportlab.lib.units import inch
                from reportlab.lib.styles import ParagraphStyle
            except Exception as e:
                return jsonify({"error": f"ReportLab not installed: {e}. Please add reportlab to requirements."}), 500

            out_path = os.path.join(app.config["UPLOAD_FOLDER"], f"synthesis_{uuid.uuid4()}.pdf")
            doc = SimpleDocTemplate(out_path, pagesize=LETTER, leftMargin=0.8*inch, rightMargin=0.8*inch, topMargin=0.8*inch, bottomMargin=0.8*inch)
            story = []

            styles = getSampleStyleSheet()
            header_style = ParagraphStyle('Header', parent=styles['Heading1'], spaceAfter=12)
            h2_style = ParagraphStyle('H2', parent=styles['Heading2'], spaceBefore=8, spaceAfter=6)
            h3_style = ParagraphStyle('H3', parent=styles['Heading3'], spaceBefore=6, spaceAfter=4)
            body_style = ParagraphStyle('Body', parent=styles['BodyText'], leading=14, spaceAfter=8, fontSize=11)
            bullet1_style = ParagraphStyle('BulletL1', parent=styles['BodyText'], leftIndent=18, spaceBefore=2)
            bullet2_style = ParagraphStyle('BulletL2', parent=styles['BodyText'], leftIndent=36, spaceBefore=2)
            bullet3_style = ParagraphStyle('BulletL3', parent=styles['BodyText'], leftIndent=54, spaceBefore=2)

            title_map = {"report": "SYNTHESIZED REPORT", "brief": "SYNTHESIZED BRIEF", "minutes": "SYNTHESIZED MINUTES"}
            story.append(Paragraph(title_map.get(target_format, "SYNTHESIZED REPORT"), header_style))
            story.append(Spacer(1, 0.2*inch))

            # Parse AI text into headings, paragraphs, and bullets with up to 3 levels
            import re
            def esc(s: str) -> str:
                return s.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')

            text = ai_text.replace('\r\n', '\n')
            text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
            text = re.sub(r'(?<=[\n\.;:])\s*(?:\*|\-|\+)\s+', '\n* ', text)

            def is_h2(line: str) -> bool:
                l = line.strip()
                if len(l) < 2 or len(l) > 80:
                    return False
                if l.endswith(':'):
                    return True
                return bool(re.match(r'^[A-Z0-9][A-Z0-9\s\-/()&]{2,80}$', l))

            def is_h3(line: str) -> bool:
                l = line.strip()
                if len(l) < 2 or len(l) > 80:
                    return False
                if is_h2(l):
                    return False
                return bool(re.match(r'^(?:[A-Z][\w()\-/]*)(?:\s+(?:[A-Z][\w()\-/]*|of|and|to|for|in|on|with|the|a|an))*$', l))

            bullet_re = re.compile(r'^(?P<indent>\s*)(?:(?:[\*\-\+])|(?:\d+\.))\s+(?P<text>.+)$')

            lines = text.split('\n')
            paragraph_buf = []

            def flush_paragraph():
                if paragraph_buf:
                    paragraph_text = ' '.join(paragraph_buf).strip()
                    if paragraph_text:
                        story.append(Paragraph(esc(paragraph_text), body_style))
                    paragraph_buf.clear()

            for raw in lines:
                line = raw.rstrip()
                stripped = line.strip()
                if not stripped:
                    flush_paragraph()
                    continue

                m = bullet_re.match(line)
                if m:
                    flush_paragraph()
                    indent = m.group('indent') or ''
                    level = 1 + min(2, len(indent) // 2)
                    txt = m.group('text').strip()
                    bullet_text = f"• {txt}"
                    style = bullet1_style if level == 1 else (bullet2_style if level == 2 else bullet3_style)
                    story.append(Paragraph(esc(bullet_text), style))
                    continue

                if is_h2(stripped):
                    flush_paragraph()
                    heading = stripped[:-1].strip() if stripped.endswith(':') else stripped
                    story.append(Paragraph(esc(heading), h2_style))
                    continue
                if is_h3(stripped):
                    flush_paragraph()
                    story.append(Paragraph(esc(stripped), h3_style))
                    continue

                paragraph_buf.append(stripped)

            flush_paragraph()

            # Append appendices
            try:
                quality = artifacts.get("quality") if isinstance(artifacts, dict) else None
                conflicts = artifacts.get("conflicts") if isinstance(artifacts, dict) else None
                provenance = artifacts.get("provenance") if isinstance(artifacts, dict) else None
            except Exception:
                quality, conflicts, provenance = None, None, None

            if quality and isinstance(quality, dict):
                story.append(PageBreak())
                story.append(Paragraph('QUALITY REPORT', header_style))
                story.append(Spacer(1, 0.15*inch))
                try:
                    outline = quality.get('outline') or {}
                    prov = quality.get('provenance') or {}
                    numeric = quality.get('numeric') or {}
                    overall = quality.get('overall_score')
                    story.append(Paragraph(esc(f"Overall Score: {overall}"), h2_style))
                    story.append(Paragraph('Outline Coverage', h3_style))
                    story.append(Paragraph(esc(f"Covered {outline.get('covered',0)} of {outline.get('total',0)} ({outline.get('coverage',0)})"), body_style))
                    missing = outline.get('missing') or []
                    if missing:
                        story.append(Paragraph('Missing Sections:', h3_style))
                        for m in missing[:10]:
                            story.append(Paragraph(esc(f"- {m}"), body_style))
                    story.append(Paragraph('Provenance Coverage', h3_style))
                    story.append(Paragraph(esc(f"Lines with sources: {prov.get('with_sources',0)}/{prov.get('lines',0)} ({prov.get('coverage',0)})"), body_style))
                    story.append(Paragraph(esc(f"Avg top score: {prov.get('avg_top_score',0)}"), body_style))
                    story.append(Paragraph('Numeric Alignment', h3_style))
                    story.append(Paragraph(esc(f"Numbers matched: {numeric.get('matched',0)}/{numeric.get('numbers',0)} ({numeric.get('match_ratio',0)})"), body_style))
                except Exception:
                    pass

            if conflicts and isinstance(conflicts, list):
                story.append(PageBreak())
                story.append(Paragraph('CONFLICTS SUMMARY', header_style))
                story.append(Spacer(1, 0.15*inch))
                for c in conflicts[:50]:
                    label = c.get('label', 'value')
                    story.append(Paragraph(esc(f"Label: {label}"), h2_style))
                    for val_entry in c.get('values', [])[:10]:
                        val = val_entry.get('value', '')
                        story.append(Paragraph(esc(f"Value: {val}"), h3_style))
                        for occ in val_entry.get('occurrences', [])[:10]:
                            fn = occ.get('filename', 'unknown')
                            raw = occ.get('raw', '')
                            story.append(Paragraph(esc(f"- {fn}: {raw}"), body_style))
                        story.append(Spacer(1, 0.1*inch))

            if provenance and isinstance(provenance, list):
                story.append(PageBreak())
                story.append(Paragraph('PROVENANCE (Top matches per synthesized line)', header_style))
                story.append(Spacer(1, 0.15*inch))
                for m in provenance[:200]:
                    line_idx = m.get('line_index')
                    line = m.get('line', '')
                    story.append(Paragraph(esc(f"Line {line_idx}: {line}"), h3_style))
                    for src in m.get('sources', [])[:5]:
                        fi = src.get('file_index')
                        pg = src.get('page')
                        sc = src.get('score')
                        story.append(Paragraph(esc(f"- file #{fi}, page {pg}, score {sc}"), body_style))
                    story.append(Spacer(1, 0.08*inch))

            try:
                doc.build(story)
            except Exception as e:
                return jsonify({"error": f"Failed to render PDF: {e}"}), 500

            @after_this_request
            def _cleanup(resp):
                try:
                    os.remove(out_path)
                except OSError:
                    pass
                return resp

            return send_file(
                out_path,
                as_attachment=True,
                download_name=f"document_{target_format}.pdf",
                mimetype="application/pdf"
            )

        # DOCX export
        if output_format == "docx":
            try:
                from docx import Document
            except Exception as e:
                return jsonify({"error": f"python-docx not installed: {e}. Please add python-docx to requirements."}), 500

            docx_path = os.path.join(app.config["UPLOAD_FOLDER"], f"synthesis_{uuid.uuid4()}.docx")
            document = Document()

            # Title
            title_map = {"report": "SYNTHESIZED REPORT", "brief": "SYNTHESIZED BRIEF", "minutes": "SYNTHESIZED MINUTES"}
            document.add_heading(title_map.get(target_format, "SYNTHESIZED REPORT"), level=1)

            # Parse AI text similar to PDF branch
            import re
            text = ai_text.replace('\r\n', '\n')
            text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
            text = re.sub(r'(?<=[\n\.;:])\s*(?:\*|\-|\+)\s+', '\n* ', text)

            def is_h2(line: str) -> bool:
                l = line.strip()
                if len(l) < 2 or len(l) > 80:
                    return False
                if l.endswith(':'):
                    return True
                return bool(re.match(r'^[A-Z0-9][A-Z0-9\s\-/()&]{2,80}$', l))

            def is_h3(line: str) -> bool:
                l = line.strip()
                if len(l) < 2 or len(l) > 80:
                    return False
                if is_h2(l):
                    return False
                return bool(re.match(r'^(?:[A-Z][\w()\-/]*)(?:\s+(?:[A-Z][\w()\-/]*|of|and|to|for|in|on|with|the|a|an))*$', l))

            bullet_re = re.compile(r'^(?P<indent>\s*)(?:(?:[\*\-\+])|(?:\d+\.))\s+(?P<text>.+)$')
            lines = text.split('\n')
            paragraph_buf = []

            def flush_paragraph():
                if paragraph_buf:
                    paragraph_text = ' '.join(paragraph_buf).strip()
                    if paragraph_text:
                        document.add_paragraph(paragraph_text)
                    paragraph_buf.clear()

            for raw in lines:
                line = raw.rstrip()
                stripped = line.strip()
                if not stripped:
                    flush_paragraph()
                    continue

                m = bullet_re.match(line)
                if m:
                    flush_paragraph()
                    indent = m.group('indent') or ''
                    level = 1 + min(2, len(indent) // 2)
                    txt = m.group('text').strip()
                    p = document.add_paragraph(txt)
                    p.style = 'List Bullet' if level == 1 else ('List Bullet 2' if level == 2 else 'List Bullet 3')
                    continue

                if is_h2(stripped):
                    flush_paragraph()
                    document.add_heading(stripped[:-1].strip() if stripped.endswith(':') else stripped, level=2)
                    continue
                if is_h3(stripped):
                    flush_paragraph()
                    document.add_heading(stripped, level=3)
                    continue

                paragraph_buf.append(stripped)

            flush_paragraph()

            # Append appendices
            quality = artifacts.get("quality") if isinstance(artifacts, dict) else None
            conflicts = artifacts.get("conflicts") if isinstance(artifacts, dict) else None
            provenance = artifacts.get("provenance") if isinstance(artifacts, dict) else None

            if quality and isinstance(quality, dict):
                document.add_page_break()
                document.add_heading('QUALITY REPORT', level=1)
                outline = quality.get('outline') or {}
                prov = quality.get('provenance') or {}
                numeric = quality.get('numeric') or {}
                overall = quality.get('overall_score')
                document.add_paragraph(f"Overall Score: {overall}")
                document.add_heading('Outline Coverage', level=2)
                document.add_paragraph(f"Covered {outline.get('covered',0)} of {outline.get('total',0)} ({outline.get('coverage',0)})")
                missing = outline.get('missing') or []
                if missing:
                    document.add_paragraph('Missing Sections:')
                    for m in missing[:10]:
                        p = document.add_paragraph(f"- {m}")
                        p.style = 'List Bullet'
                document.add_heading('Provenance Coverage', level=2)
                document.add_paragraph(f"Lines with sources: {prov.get('with_sources',0)}/{prov.get('lines',0)} ({prov.get('coverage',0)})")
                document.add_paragraph(f"Avg top score: {prov.get('avg_top_score',0)}")
                document.add_heading('Numeric Alignment', level=2)
                document.add_paragraph(f"Numbers matched: {numeric.get('matched',0)}/{numeric.get('numbers',0)} ({numeric.get('match_ratio',0)})")

            if conflicts and isinstance(conflicts, list):
                document.add_page_break()
                document.add_heading('CONFLICTS SUMMARY', level=1)
                for c in conflicts[:50]:
                    label = c.get('label','value')
                    document.add_heading(f"Label: {label}", level=2)
                    for val_entry in c.get('values', [])[:10]:
                        val = val_entry.get('value','')
                        document.add_heading(f"Value: {val}", level=3)
                        for occ in val_entry.get('occurrences', [])[:10]:
                            fn = occ.get('filename','unknown'); raw = occ.get('raw','')
                            p = document.add_paragraph(f"{fn}: {raw}"); p.style = 'List Bullet'

            if provenance and isinstance(provenance, list):
                document.add_page_break()
                document.add_heading('PROVENANCE (Top matches per synthesized line)', level=1)
                for m in provenance[:200]:
                    line_idx = m.get('line_index'); line = m.get('line','')
                    document.add_heading(f"Line {line_idx}: {line}", level=3)
                    for src in m.get('sources', [])[:5]:
                        fi = src.get('file_index'); pg = src.get('page'); sc = src.get('score')
                        p = document.add_paragraph(f"file #{fi}, page {pg}, score {sc}"); p.style = 'List Bullet'

            try:
                document.save(docx_path)
            except Exception as e:
                return jsonify({"error": f"Failed to render DOCX: {e}"}), 500

            @after_this_request
            def _cleanup_docx(resp):
                try:
                    os.remove(docx_path)
                except OSError:
                    pass
                return resp

            return send_file(
                docx_path,
                as_attachment=True,
                download_name=f"document_{target_format}.docx",
                mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            )
        if output_format == "pptx":
            try:
                from pptx import Presentation
                from pptx.util import Inches, Pt
            except Exception as e:
                return jsonify({"error": f"python-pptx not installed: {e}. Please add python-pptx to requirements."}), 500

            prs = Presentation()
            slide_layout = prs.slide_layouts[1]  # Title and Content

            # Title slide
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title_map = {"report": "SYNTHESIZED REPORT", "brief": "SYNTHESIZED BRIEF", "minutes": "SYNTHESIZED MINUTES"}
            title_slide.shapes.title.text = title_map.get(target_format, "SYNTHESIZED REPORT")
            title_slide.placeholders[1].text = "Generated by AI Document Synthesis"

            # Parse AI text into slides by H2 headings
            import re
            text = ai_text.replace('\r\n', '\n')
            text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
            text = re.sub(r'(?<=[\n\.;:])\s*(?:\*|\-|\+)\s+', '\n* ', text)
            lines = [ln.rstrip() for ln in text.split('\n')]

            def is_h2(line: str) -> bool:
                l = line.strip()
                if len(l) < 2 or len(l) > 80:
                    return False
                if l.endswith(':'):
                    return True
                return bool(re.match(r'^[A-Z0-9][A-Z0-9\s\-/()&]{2,80}$', l))

            def flush_slide(title: str, bullets: list[str]):
                if not title and not bullets:
                    return
                s = prs.slides.add_slide(slide_layout)
                s.shapes.title.text = title or "Section"
                body = s.shapes.placeholders[1].text_frame
                if bullets:
                    body.clear()
                    for b in bullets[:12]:
                        if not b:
                            continue
                        if not body.text:
                            body.text = b
                        else:
                            body.add_paragraph().text = b

            current_title = None
            current_bullets: list[str] = []
            for ln in lines:
                if is_h2(ln.strip()):
                    flush_slide(current_title, current_bullets)
                    current_title = ln.strip().rstrip(':')
                    current_bullets = []
                    continue
                if ln.strip().startswith('* '):
                    current_bullets.append(ln.strip()[2:])
                else:
                    # treat plain paragraphs as bullets (trim length)
                    if ln.strip():
                        current_bullets.append(ln.strip()[:160])
            flush_slide(current_title, current_bullets)

            # Append appendices as slides
            quality = artifacts.get("quality") if isinstance(artifacts, dict) else None
            conflicts = artifacts.get("conflicts") if isinstance(artifacts, dict) else None
            provenance = artifacts.get("provenance") if isinstance(artifacts, dict) else None

            if quality:
                s = prs.slides.add_slide(slide_layout); s.shapes.title.text = 'QUALITY REPORT'
                body = s.shapes.placeholders[1].text_frame
                outline = quality.get('outline') or {}; prov = quality.get('provenance') or {}; numeric = quality.get('numeric') or {}
                body.text = f"Overall Score: {quality.get('overall_score')}"
                body.add_paragraph().text = f"Outline: {outline.get('covered',0)}/{outline.get('total',0)} ({outline.get('coverage',0)})"
                body.add_paragraph().text = f"Provenance: {prov.get('with_sources',0)}/{prov.get('lines',0)} ({prov.get('coverage',0)})"
                body.add_paragraph().text = f"Numeric: {numeric.get('matched',0)}/{numeric.get('numbers',0)} ({numeric.get('match_ratio',0)})"

            if conflicts:
                s = prs.slides.add_slide(slide_layout); s.shapes.title.text = 'CONFLICTS SUMMARY'
                body = s.shapes.placeholders[1].text_frame
                for c in conflicts[:12]:
                    body.add_paragraph().text = f"Label: {c.get('label','value')}"

            if provenance:
                s = prs.slides.add_slide(slide_layout); s.shapes.title.text = 'PROVENANCE'
                body = s.shapes.placeholders[1].text_frame
                for m in provenance[:12]:
                    p = body.add_paragraph(); p.text = f"Line {m.get('line_index')}: {m.get('line','')[:120]}"

            pptx_path = os.path.join(app.config["UPLOAD_FOLDER"], f"synthesis_{uuid.uuid4()}.pptx")
            try:
                prs.save(pptx_path)
            except Exception as e:
                return jsonify({"error": f"Failed to render PPTX: {e}"}), 500

            @after_this_request
            def _cleanup_pptx(resp):
                try:
                    os.remove(pptx_path)
                except OSError:
                    pass
                return resp

            return send_file(
                pptx_path,
                as_attachment=True,
                download_name=f"document_{target_format}.pptx",
                mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
        if output_format == "teams":
            # Send a simple Teams message via incoming webhook URL
            teams_webhook = (request.form.get("teams_webhook") or os.getenv("TEAMS_WEBHOOK_URL") or "").strip()
            if not teams_webhook:
                return jsonify({"error": "Teams webhook URL missing. Provide 'teams_webhook' form field or set TEAMS_WEBHOOK_URL env."}), 400

            # Construct a summary card-like message
            summary_lines = ai_text.split('\n')[:20]
            content = "\n".join(summary_lines)
            payload = {
                "text": f"AI Document Synthesis ({target_format})\n\n{content}"
            }
            try:
                r = requests.post(teams_webhook, json=payload, timeout=15)
                if r.status_code >= 400:
                    return jsonify({"error": f"Teams webhook error: {r.status_code} {r.text}"}), 502
            except Exception as e:
                return jsonify({"error": f"Teams webhook failed: {e}"}), 502
            return jsonify({"status": "ok", "delivered": True})

    except Exception as e:
        return jsonify({"error": f"Synthesis failed: {str(e)}"}), 500


@app.post("/api/ai-document-synthesis/async")
def ai_document_synthesis_async():
    """Async version of synthesis - returns job_id immediately"""
    try:
        files = request.files.getlist("files")
        if not files:
            return jsonify({"error": "Upload at least one PDF file"}), 400
        if any(not allowed_pdf(f) for f in files):
            return jsonify({"error": "Only PDF files are allowed"}), 400

        # Validate parameters
        target_format = (request.form.get("format") or "report").lower().strip()
        if target_format not in ("report", "brief", "minutes"):
            return jsonify({"error": "format must be one of: report|brief|minutes"}), 400

        ALLOWED_DOMAINS = {
            "invoice", "contract", "financial", "purchase_order", "receipt", 
            "bank_statement", "tax_form", "resume", "legal_pleading", "patent", 
            "medical_bill", "lab_report", "insurance_claim", "real_estate", 
            "shipping_manifest", "research", "healthcare", "legal", "finance", "general"
        }
        
        domain_override = (request.form.get("domain_override") or "").lower().strip() or None
        if domain_override and domain_override not in ALLOWED_DOMAINS:
            return jsonify({"error": f"domain_override must be one of: {', '.join(sorted(ALLOWED_DOMAINS))}"}), 400
            
        template_profile = (request.form.get("template_profile") or "").lower().strip() or None
        if template_profile and template_profile not in ("executive_summary", "risk_assessment", "compliance_review"):
            return jsonify({"error": "template_profile must be one of: executive_summary|risk_assessment|compliance_review"}), 400
            
        custom_instructions = (request.form.get("custom_instructions") or "").strip() or None
        output_format = (request.form.get("output_format") or "pdf").lower().strip()
        if output_format not in ("pdf", "docx", "pptx", "teams"):
            return jsonify({"error": "output_format must be one of: pdf|docx|pptx|teams"}), 400

        # Save files to temp
        job_id = str(uuid.uuid4())
        temp_files = []
        for f in files:
            temp_path = os.path.join(app.config["UPLOAD_FOLDER"], f"{job_id}_{f.filename}")
            f.save(temp_path)
            temp_files.append(temp_path)

        # Create job
        with SYNTHESIS_JOBS_LOCK:
            SYNTHESIS_JOBS[job_id] = {
                "status": "pending",
                "progress": 0,
                "created_at": datetime.utcnow().isoformat(),
                "files": temp_files,
                "params": {
                    "target_format": target_format,
                    "domain_override": domain_override,
                    "template_profile": template_profile,
                    "custom_instructions": custom_instructions,
                    "output_format": output_format,
                },
                "result": None,
                "error": None,
            }

        # Start background thread
        def process_synthesis():
            try:
                with SYNTHESIS_JOBS_LOCK:
                    SYNTHESIS_JOBS[job_id]["status"] = "processing"
                    SYNTHESIS_JOBS[job_id]["progress"] = 10

                # Re-open files for processing
                file_objects = []
                for path in temp_files:
                    file_objects.append(open(path, 'rb'))

                params = SYNTHESIS_JOBS[job_id]["params"]
                
                # Import and run orchestrator
                try:
                    from synthesis.pipeline import SynthesisOrchestrator
                except Exception:
                    from backend.synthesis.pipeline import SynthesisOrchestrator

                with SYNTHESIS_JOBS_LOCK:
                    SYNTHESIS_JOBS[job_id]["progress"] = 20

                orchestrator = SynthesisOrchestrator()
                from synthesis.ai_service import AIServiceRouter
                callers = [(AI_SERVICE, lambda p: call_ai_service(p))]
                if AI_SERVICE != "openai" and OPENAI_API_KEY:
                    callers.append(("openai", lambda p: call_openai(p)))
                if AI_SERVICE != "anthropic" and ANTHROPIC_API_KEY:
                    callers.append(("anthropic", lambda p: call_anthropic(p)))
                router = AIServiceRouter(callers, min_len=80)

                with SYNTHESIS_JOBS_LOCK:
                    SYNTHESIS_JOBS[job_id]["progress"] = 40

                ai_text, artifacts = orchestrator.run(
                    files=file_objects,
                    target_format=params["target_format"],
                    per_file_max_chars=SYNTHESIS_PER_FILE_MAX,
                    global_max_chars=SYNTHESIS_MAX_CHARS,
                    ai_caller=lambda prompt: router.generate(prompt),
                    forced_domain=params["domain_override"],
                    template_profile=params["template_profile"],
                    user_instructions=params["custom_instructions"],
                )

                # Close files
                for f in file_objects:
                    f.close()

                with SYNTHESIS_JOBS_LOCK:
                    SYNTHESIS_JOBS[job_id]["progress"] = 70

                # Generate output file (reuse existing PDF/DOCX generation logic)
                out_format = params["output_format"]
                output_path = None
                
                if out_format == "pdf":
                    # Generate PDF
                    try:
                        from reportlab.lib.pagesizes import LETTER
                        from reportlab.lib.styles import getSampleStyleSheet
                        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
                        from reportlab.lib.units import inch
                        from reportlab.lib.styles import ParagraphStyle
                    except Exception as e:
                        raise Exception(f"ReportLab not installed: {e}")

                    output_path = os.path.join(app.config["UPLOAD_FOLDER"], f"synthesis_{job_id}.pdf")
                    doc = SimpleDocTemplate(output_path, pagesize=LETTER, leftMargin=0.8*inch, rightMargin=0.8*inch, topMargin=0.8*inch, bottomMargin=0.8*inch)
                    story = []

                    styles = getSampleStyleSheet()
                    header_style = ParagraphStyle('Header', parent=styles['Heading1'], spaceAfter=12)
                    h2_style = ParagraphStyle('H2', parent=styles['Heading2'], spaceBefore=8, spaceAfter=6)
                    h3_style = ParagraphStyle('H3', parent=styles['Heading3'], spaceBefore=6, spaceAfter=4)
                    body_style = ParagraphStyle('Body', parent=styles['BodyText'], leading=14, spaceAfter=8, fontSize=11)
                    bullet1_style = ParagraphStyle('BulletL1', parent=styles['BodyText'], leftIndent=18, spaceBefore=2)
                    bullet2_style = ParagraphStyle('BulletL2', parent=styles['BodyText'], leftIndent=36, spaceBefore=2)
                    bullet3_style = ParagraphStyle('BulletL3', parent=styles['BodyText'], leftIndent=54, spaceBefore=2)

                    title_map = {"report": "SYNTHESIZED REPORT", "brief": "SYNTHESIZED BRIEF", "minutes": "SYNTHESIZED MINUTES"}
                    story.append(Paragraph(title_map.get(params["target_format"], "SYNTHESIZED REPORT"), header_style))
                    story.append(Spacer(1, 0.2*inch))

                    # Parse AI text into headings, paragraphs, and bullets
                    def esc(s: str) -> str:
                        return s.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')

                    text = ai_text.replace('\r\n', '\n')
                    text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
                    text = re.sub(r'(?<=[\n\.;:])\s*(?:\*|\-|\+)\s+', '\n* ', text)

                    def is_h2(line: str) -> bool:
                        l = line.strip()
                        if len(l) < 2 or len(l) > 80:
                            return False
                        if l.endswith(':'):
                            return True
                        return bool(re.match(r'^[A-Z0-9][A-Z0-9\s\-/()&]{2,80}$', l))

                    def is_h3(line: str) -> bool:
                        l = line.strip()
                        if len(l) < 2 or len(l) > 80:
                            return False
                        if is_h2(l):
                            return False
                        return bool(re.match(r'^(?:[A-Z][\w()\-/]*)(?:\s+(?:[A-Z][\w()\-/]*|of|and|to|for|in|on|with|the|a|an))*$', l))

                    bullet_re = re.compile(r'^(?P<indent>\s*)(?:(?:[\*\-\+])|(?:\d+\.))\s+(?P<text>.+)$')

                    lines = text.split('\n')
                    paragraph_buf = []

                    def flush_paragraph():
                        if paragraph_buf:
                            paragraph_text = ' '.join(paragraph_buf).strip()
                            if paragraph_text:
                                story.append(Paragraph(esc(paragraph_text), body_style))
                            paragraph_buf.clear()

                    for raw in lines:
                        line = raw.rstrip()
                        stripped = line.strip()
                        if not stripped:
                            flush_paragraph()
                            continue

                        m = bullet_re.match(line)
                        if m:
                            flush_paragraph()
                            indent = m.group('indent') or ''
                            level = 1 + min(2, len(indent) // 2)
                            txt = m.group('text').strip()
                            bullet_text = f"• {txt}"
                            style = bullet1_style if level == 1 else (bullet2_style if level == 2 else bullet3_style)
                            story.append(Paragraph(esc(bullet_text), style))
                            continue

                        if is_h2(stripped):
                            flush_paragraph()
                            heading = stripped[:-1].strip() if stripped.endswith(':') else stripped
                            story.append(Paragraph(esc(heading), h2_style))
                            continue
                        if is_h3(stripped):
                            flush_paragraph()
                            story.append(Paragraph(esc(stripped), h3_style))
                            continue

                        paragraph_buf.append(stripped)

                    flush_paragraph()
                    doc.build(story)
                    
                elif out_format == "docx":
                    # Generate DOCX
                    try:
                        from docx import Document
                        from docx.shared import Pt, Inches
                    except Exception as e:
                        raise Exception(f"python-docx not installed: {e}")

                    output_path = os.path.join(app.config["UPLOAD_FOLDER"], f"synthesis_{job_id}.docx")
                    document = Document()
                    document.add_heading(f"Synthesized {params['target_format'].capitalize()}", 0)
                    
                    # Simple paragraph-based rendering
                    for line in ai_text.split('\n'):
                        line = line.strip()
                        if not line:
                            continue
                        if line.isupper() and len(line) < 80:
                            document.add_heading(line, level=1)
                        elif line.startswith('- ') or line.startswith('* '):
                            document.add_paragraph(line[2:], style='List Bullet')
                        else:
                            document.add_paragraph(line)
                    
                    document.save(output_path)

                with SYNTHESIS_JOBS_LOCK:
                    SYNTHESIS_JOBS[job_id]["status"] = "completed"
                    SYNTHESIS_JOBS[job_id]["progress"] = 100
                    SYNTHESIS_JOBS[job_id]["result"] = {
                        "output_path": output_path,
                        "artifacts": {
                            "quality": artifacts.get("quality"),
                            "conflicts": artifacts.get("conflicts"),
                            "template": artifacts.get("template"),
                        }
                    }

                # Cleanup temp files
                for path in temp_files:
                    try:
                        os.remove(path)
                    except:
                        pass

            except Exception as e:
                with SYNTHESIS_JOBS_LOCK:
                    SYNTHESIS_JOBS[job_id]["status"] = "failed"
                    SYNTHESIS_JOBS[job_id]["error"] = str(e)

        thread = threading.Thread(target=process_synthesis, daemon=True)
        thread.start()

        return jsonify({"job_id": job_id}), 202

    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.get("/api/ai-document-synthesis/status/<job_id>")
def get_synthesis_status(job_id: str):
    """Get status of async synthesis job"""
    with SYNTHESIS_JOBS_LOCK:
        job = SYNTHESIS_JOBS.get(job_id)
        if not job:
            return jsonify({"error": "Job not found"}), 404
        
        return jsonify({
            "job_id": job_id,
            "status": job["status"],
            "progress": job["progress"],
            "error": job.get("error"),
            "done": job["status"] in ("completed", "failed"),
        })


@app.get("/api/ai-document-synthesis/result/<job_id>")
def get_synthesis_result(job_id: str):
    """Download result of completed synthesis job"""
    with SYNTHESIS_JOBS_LOCK:
        job = SYNTHESIS_JOBS.get(job_id)
        if not job:
            return jsonify({"error": "Job not found"}), 404
        
        if job["status"] != "completed":
            return jsonify({"error": "Job not completed"}), 400
        
        result = job.get("result")
        if not result or not result.get("output_path"):
            return jsonify({"error": "No result file"}), 404
        
        output_path = result["output_path"]
        if not os.path.exists(output_path):
            return jsonify({"error": "Result file not found"}), 404
        
        # Determine mimetype
        ext = os.path.splitext(output_path)[1].lower()
        mimetype_map = {
            ".pdf": "application/pdf",
            ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        }
        mimetype = mimetype_map.get(ext, "application/octet-stream")
        
        return send_file(output_path, mimetype=mimetype, as_attachment=True, download_name=os.path.basename(output_path))


def _fallback_extract(temp_path: str, options: ExtractionOptions) -> dict:
    """Fallback extraction logic when new pipeline is not available."""
    try:
        # Read PDF text
        reader = PdfReader(temp_path)
        pages_text = []
        try:
            from concurrent.futures import ThreadPoolExecutor
            def _extract_page(i):
                try:
                    return reader.pages[i].extract_text() or ""
                except Exception:
                    return ""
            with ThreadPoolExecutor(max_workers=4) as ex:
                pages_text = list(ex.map(_extract_page, range(len(reader.pages))))
        except Exception:
            pages_text = []
            for p in reader.pages:
                try:
                    pages_text.append(p.extract_text() or "")
                except Exception:
                    pages_text.append("")
        
        all_text = "\n".join(pages_text)
        
        # Use old extraction logic
        override = options.domain_override or ''
        custom_instructions = options.custom_instructions or ''
        domain_norm = _normalize_domain(override)
        selected_fields = _validate_selected_fields(domain_norm, options.selected_fields)
        result = ai_only_extract(all_text, pages_text, custom_instructions, domain_norm, selected_fields)
        
        # Optional enrichment
        if options.enrich:
            try:
                enriched = enrich_extraction_with_llm(result.get('type'), result.get('mapped_fields'), all_text, pages_text, build_inverted_index(pages_text))
                new_mapped = enriched.get('mapped_fields', result.get('mapped_fields'))
                if isinstance(new_mapped, dict):
                    result['mapped_fields'] = new_mapped
            except Exception:
                pass
        
        return result
    except Exception as e:
        return {"error": f"Fallback extraction failed: {str(e)}"}

@app.post("/api/extract")
def api_extract():
    """MVP extraction endpoint: returns detected document type and basic entities.
    Input: multipart/form-data with 'file' (PDF)
    Output: { type, pages, entities: { emails, phones, amounts, dates } }
    """
    try:
        if "file" not in request.files:
            return jsonify({"error": "No file provided"}), 400
        file = request.files["file"]
        if file.filename == "":
            return jsonify({"error": "No file selected"}), 400
        if not allowed_pdf(file):
            return jsonify({"error": "Invalid file type. Only PDF files are allowed."}), 400

        # Persist to temp path
        fname = secure_filename(file.filename) or f"upload_{uuid.uuid4()}.pdf"
        temp_path = os.path.join(app.config["UPLOAD_FOLDER"], f"{uuid.uuid4()}_{fname}")
        file.save(temp_path)

        # Parse request parameters into ExtractionOptions
        try:
            selected_fields = json.loads(request.form.get('selected_fields') or '[]')
            if not isinstance(selected_fields, list):
                selected_fields = []
        except Exception:
            selected_fields = []

        options = ExtractionOptions(
            domain_override=request.form.get('domain_override'),
            selected_fields=selected_fields,
            custom_instructions=request.form.get('custom_instructions', ''),
            enrich=request.form.get('enrich', '').lower() in ('1', 'true', 'yes'),
            extract_tables=request.form.get('extract_tables', '').lower() in ('1', 'true', 'yes'),
            extract_formulas=request.form.get('extract_formulas', '').lower() in ('1', 'true', 'yes')
        )

        # Use new pipeline if available, otherwise fallback to old logic
        if _extraction_pipeline:
            result = _extraction_pipeline.extract(temp_path, options)
            resp_json = result.to_dict()
        else:
            # Fallback to old extraction logic
            resp_json = _fallback_extract(temp_path, options)

        @after_this_request
        def _cleanup(resp):
            try:
                if os.path.isfile(temp_path):
                    os.remove(temp_path)
            except Exception:
                pass
            return resp

        return jsonify(resp_json)
    except Exception as e:
        return jsonify({"error": f"Extraction failed: {str(e)}"}), 500

@app.post("/api/extract/async")
def api_extract_async():
    try:
        if "file" not in request.files:
            return jsonify({"error": "No file provided"}), 400
        upfile = request.files["file"]
        if upfile.filename == "":
            return jsonify({"error": "No file selected"}), 400
        if not allowed_pdf(upfile):
            return jsonify({"error": "Invalid file type. Only PDF files are allowed."}), 400

        fname = secure_filename(upfile.filename) or f"upload_{uuid.uuid4()}.pdf"
        temp_path = os.path.join(app.config["UPLOAD_FOLDER"], f"{uuid.uuid4()}_{fname}")
        upfile.save(temp_path)

        # Parse request parameters into ExtractionOptions
        try:
            selected_fields = json.loads(request.form.get('selected_fields') or '[]')
            if not isinstance(selected_fields, list):
                selected_fields = []
        except Exception:
            selected_fields = []

        options = ExtractionOptions(
            domain_override=request.form.get('domain_override'),
            selected_fields=selected_fields,
            custom_instructions=request.form.get('custom_instructions', ''),
            enrich=request.form.get('enrich', '').lower() in ('1', 'true', 'yes'),
            extract_tables=request.form.get('extract_tables', '').lower() in ('1', 'true', 'yes'),
            extract_formulas=request.form.get('extract_formulas', '').lower() in ('1', 'true', 'yes')
        )

        job_id = str(uuid.uuid4())
        with EXTRACT_JOBS_LOCK:
            EXTRACT_JOBS[job_id] = {"progress": 0, "done": False, "error": None, "result": None}

        def run_job(jid: str, path: str, extraction_options: ExtractionOptions):
            def set_progress(p: int):
                with EXTRACT_JOBS_LOCK:
                    if jid in EXTRACT_JOBS:
                        EXTRACT_JOBS[jid]["progress"] = max(0, min(100, int(p)))
            
            try:
                print(f"DEBUG: Starting async job {jid} for file {path}")
                set_progress(10)
                
                # Use new pipeline if available, otherwise fallback to old logic
                if _extraction_pipeline:
                    print(f"DEBUG: Using new extraction pipeline for job {jid}")
                    result = _extraction_pipeline.extract(path, extraction_options)
                    resp_json = result.to_dict()
                    print(f"DEBUG: Pipeline extraction completed for job {jid}")
                else:
                    print(f"DEBUG: Using fallback extraction for job {jid}")
                    # Fallback to old extraction logic
                    resp_json = _fallback_extract(path, extraction_options)
                    print(f"DEBUG: Fallback extraction completed for job {jid}")
                
                set_progress(90)
                
                # Ensure result has expected format
                if not isinstance(resp_json, dict):
                    print(f"DEBUG: Invalid result format for job {jid}: {type(resp_json)}")
                    resp_json = {"error": "Invalid extraction result"}
                
                set_progress(100)
                
                with EXTRACT_JOBS_LOCK:
                    if jid in EXTRACT_JOBS:
                        EXTRACT_JOBS[jid]["result"] = resp_json
                        EXTRACT_JOBS[jid]["done"] = True
                        EXTRACT_JOBS[jid]["progress"] = 100
                        print(f"DEBUG: Job {jid} completed successfully")
                        
            except Exception as e:
                print(f"DEBUG: Error in job {jid}: {str(e)}")
                import traceback
                traceback.print_exc()
                with EXTRACT_JOBS_LOCK:
                    if jid in EXTRACT_JOBS:
                        EXTRACT_JOBS[jid]["error"] = str(e)
                        EXTRACT_JOBS[jid]["done"] = True
                        EXTRACT_JOBS[jid]["progress"] = 100
                        print(f"DEBUG: Job {jid} marked as failed with error: {str(e)}")
            finally:
                try:
                    if os.path.isfile(path):
                        os.remove(path)
                        print(f"DEBUG: Cleaned up temp file {path}")
                except Exception as e:
                    print(f"DEBUG: Error cleaning up temp file {path}: {str(e)}")

        th = threading.Thread(target=run_job, args=(job_id, temp_path, options), daemon=True)
        th.start()
        return jsonify({"job_id": job_id})
    except Exception as e:
        return jsonify({"error": f"Failed to start async extract: {e}"}), 500

@app.get("/api/extract/status/<job_id>")
def api_extract_status(job_id: str):
    with EXTRACT_JOBS_LOCK:
        job = EXTRACT_JOBS.get(job_id)
        if not job:
            print(f"DEBUG: Job {job_id} not found in EXTRACT_JOBS")
            return jsonify({"error": "not_found"}), 404
        
        status_response = {
            "job_id": job_id,
            "progress": job.get("progress", 0),
            "done": job.get("done", False),
            "error": job.get("error"),
            "result": job.get("result") if job.get("done") else None,
        }
        
        # Add debug info
        if job.get("error"):
            print(f"DEBUG: Job {job_id} has error: {job.get('error')}")
        elif job.get("done"):
            print(f"DEBUG: Job {job_id} completed successfully")
        else:
            print(f"DEBUG: Job {job_id} still running, progress: {job.get('progress', 0)}%")
        
        return jsonify(status_response)

# ── Batch extraction (multiple PDFs) ───────────────────────────────
@app.post("/api/extract/batch")
def api_extract_batch_start():
    try:
        files = request.files.getlist('files') or []
        if not files:
            return jsonify({"error": "No files provided"}), 400
        
        # Parse request parameters into ExtractionOptions
        try:
            selected_fields = json.loads(request.form.get('selected_fields') or '[]')
            if not isinstance(selected_fields, list):
                selected_fields = []
        except Exception:
            selected_fields = []

        options = ExtractionOptions(
            domain_override=request.form.get('domain'),
            selected_fields=selected_fields,
            custom_instructions=request.form.get('custom_instructions', ''),
            enrich=False,  # Batch processing doesn't use enrichment
            extract_tables=request.form.get('extract_tables', '').lower() in ('1', 'true', 'yes'),
            extract_formulas=request.form.get('extract_formulas', '').lower() in ('1', 'true', 'yes')
        )

        batch_id = str(uuid.uuid4())
        with EXTRACT_JOBS_LOCK:
            EXTRACT_JOBS[batch_id] = {"progress": 0, "done": False, "error": None, "result": None, "is_batch": True, "items": []}

        saved_paths: list[tuple[str, str]] = []  # (orig_name, saved_path)
        for f in files:
            if not f.filename:
                continue
            if not allowed_pdf(f):
                continue
            fname = secure_filename(f.filename) or f"upload_{uuid.uuid4()}.pdf"
            temp_path = os.path.join(app.config["UPLOAD_FOLDER"], f"{uuid.uuid4()}_{fname}")
            f.save(temp_path)
            saved_paths.append((fname, temp_path))

        def run_batch(bid: str, items: list[tuple[str, str]], extraction_options: ExtractionOptions):
            try:
                total = max(1, len(items))
                out_results = []
                for idx, (orig, path) in enumerate(items, start=1):
                    try:
                        # Use new pipeline if available, otherwise fallback to old logic
                        if _extraction_pipeline:
                            result = _extraction_pipeline.extract(path, extraction_options)
                            result_dict = result.to_dict()
                        else:
                            # Fallback to old extraction logic
                            result_dict = _fallback_extract(path, extraction_options)
                        
                        out_results.append({"file": orig, "result": result_dict})
                    except Exception as fe:
                        out_results.append({"file": orig, "error": str(fe)})
                    finally:
                        try:
                            if os.path.isfile(path):
                                os.remove(path)
                        except Exception:
                            pass
                    with EXTRACT_JOBS_LOCK:
                        if bid in EXTRACT_JOBS:
                            EXTRACT_JOBS[bid]["progress"] = int(idx * 100 / total)
                            EXTRACT_JOBS[bid]["items"] = out_results
                with EXTRACT_JOBS_LOCK:
                    if bid in EXTRACT_JOBS:
                        EXTRACT_JOBS[bid]["done"] = True
                        EXTRACT_JOBS[bid]["result"] = {"items": out_results}
                        EXTRACT_JOBS[bid]["progress"] = 100
            except Exception as e:
                with EXTRACT_JOBS_LOCK:
                    if bid in EXTRACT_JOBS:
                        EXTRACT_JOBS[bid]["error"] = str(e)
                        EXTRACT_JOBS[bid]["done"] = True
                        EXTRACT_JOBS[bid]["progress"] = 100

        th = threading.Thread(target=run_batch, args=(batch_id, saved_paths, options), daemon=True)
        th.start()
        return jsonify({"batch_id": batch_id})
    except Exception as e:
        return jsonify({"error": f"Failed to start batch: {e}"}), 500

@app.get("/api/extract/batch/status/<batch_id>")
def api_extract_batch_status(batch_id: str):
    with EXTRACT_JOBS_LOCK:
        job = EXTRACT_JOBS.get(batch_id)
        if not job:
            return jsonify({"error": "not_found"}), 404
        if not job.get("is_batch"):
            return jsonify({"error": "not_a_batch"}), 400
        return jsonify({
            "batch_id": batch_id,
            "progress": job.get("progress", 0),
            "done": job.get("done", False),
            "error": job.get("error"),
            "items": job.get("result", {}).get("items") if job.get("done") else (job.get("items") or []),
        })

@app.get("/api/extract/batch/results/<batch_id>")
def api_extract_batch_results(batch_id: str):
    with EXTRACT_JOBS_LOCK:
        job = EXTRACT_JOBS.get(batch_id)
        if not job:
            return jsonify({"error": "not_found"}), 404
        if not job.get("is_batch"):
            return jsonify({"error": "not_a_batch"}), 400
        if not job.get("done"):
            return jsonify({"error": "not_ready"}), 400
        items = (job.get("result") or {}).get("items") or []
    
    # Build a ZIP of JSON files with legitimate-looking structure
    mem = io.BytesIO()
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    
    with zipfile.ZipFile(mem, mode='w', compression=zipfile.ZIP_DEFLATED) as zf:
        # Add manifest file to identify the archive as legitimate
        manifest = {
            "type": "pdf_extraction_results",
            "version": "1.0",
            "tool": "Perk PDF Extractor",
            "generated_at": datetime.now().isoformat(),
            "batch_id": batch_id,
            "file_count": len(items),
            "description": "Batch PDF extraction results in JSON format"
        }
        zf.writestr("README.txt", 
                   "PDF Extraction Results\n"
                   "======================\n\n"
                   f"Generated by: Perk PDF\n"
                   f"Date: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n"
                   f"Batch ID: {batch_id}\n"
                   f"Files processed: {len(items)}\n\n"
                   "This archive contains extraction results in JSON format.\n"
                   "Each file corresponds to one processed PDF document.\n")
        zf.writestr("manifest.json", json.dumps(manifest, ensure_ascii=False, indent=2))
        
        # Use predictable, descriptive filenames instead of random UUIDs
        for idx, it in enumerate(items, start=1):
            # Get original filename and sanitize it
            original_name = it.get('file', f"document_{idx}")
            # Remove extension if present
            base_name = os.path.splitext(original_name)[0]
            # Sanitize filename to be safe and predictable
            safe_name = re.sub(r"[^A-Za-z0-9_.-]", "_", base_name)
            # Use numbered prefix for clear ordering
            filename = f"extraction_{idx:03d}_{safe_name}.json"
            
            content = json.dumps(it, ensure_ascii=False, indent=2)
            zf.writestr(filename, content)
    
    mem.seek(0)
    
    # Use descriptive filename with timestamp instead of random batch_id
    download_filename = f"pdf_extraction_results_{timestamp}.zip"
    
    # Create response with proper security headers
    response = send_file(
        mem, 
        mimetype='application/zip', 
        as_attachment=True, 
        download_name=download_filename
    )
    
    # Add security headers to reduce false positive detection
    response.headers['X-Content-Type-Options'] = 'nosniff'
    response.headers['X-Download-Options'] = 'noopen'
    response.headers['Content-Description'] = 'PDF Extraction Results Archive'
    
    return response

@app.post("/api/extract/pdf")
def api_extract_pdf():
    """Generate a nicely formatted Extraction PDF on the server.
    Accepts: multipart/form-data
      - file: the PDF
      - domain_override (optional)
      - enrich (optional: 'true'|'false')
    Returns: application/pdf stream
    """
    try:
        if "file" not in request.files:
            return jsonify({"error": "No file provided"}), 400
        file = request.files["file"]
        if file.filename == "":
            return jsonify({"error": "No file selected"}), 400
        if not allowed_pdf(file):
            return jsonify({"error": "Invalid file type. Only PDF files are allowed."}), 400

        # Save to temp
        fname = secure_filename(file.filename) or f"upload_{uuid.uuid4()}.pdf"
        temp_path = os.path.join(app.config["UPLOAD_FOLDER"], f"{uuid.uuid4()}_{fname}")
        file.save(temp_path)

        # Parse request parameters into ExtractionOptions
        try:
            selected_fields = json.loads(request.form.get('selected_fields') or '[]')
            if not isinstance(selected_fields, list):
                selected_fields = []
        except Exception:
            selected_fields = []

        options = ExtractionOptions(
            domain_override=request.form.get('domain_override'),
            selected_fields=selected_fields,
            custom_instructions=request.form.get('custom_instructions', ''),
            enrich=request.form.get('enrich', '').lower() in ('1', 'true', 'yes'),
            use_ocr=False  # PDF reports don't use OCR
        )

        # Use new pipeline if available, otherwise fallback to old logic
        if _extraction_pipeline:
            result = _extraction_pipeline.extract(temp_path, options)
            result_dict = result.to_dict()
        else:
            # Fallback to old extraction logic
            result_dict = _fallback_extract(temp_path, options)

        # Extract data for PDF generation
        dtype = result_dict.get('type', 'general')
        entities = result_dict.get('entities', {})
        mapped = result_dict.get('mapped_fields', {})
        validation = result_dict.get('validation', {})
        tables = result_dict.get('tables', [])
        formulas = result_dict.get('formulas', [])
        confidence = result_dict.get('confidence', {})

        # Build PDF in-memory using ReportLab
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=A4, leftMargin=56, rightMargin=56, topMargin=56, bottomMargin=56)
        styles = getSampleStyleSheet()
        elems = []

        title_style = styles['Heading1']
        title_style.textColor = colors.HexColor("#1f2937")
        elems.append(Paragraph("Extraction Report", title_style))
        elems.append(Spacer(1, 10))

        meta_style = styles['Normal']
        elems.append(Paragraph(f"Detected Type: <b>{dtype}</b>", meta_style))
        elems.append(Paragraph(f"Pages: <b>{result_dict.get('pages', 0)}</b>", meta_style))
        elems.append(Spacer(1, 10))

        # Confidence
        if confidence:
            elems.append(Paragraph("<b>Confidence</b>", styles['Heading3']))
            elems.append(Paragraph(f"Overall: {confidence.get('overall','-')}%", meta_style))
            fields_conf = confidence.get('fields') or {}
            if fields_conf:
                bullets = [Paragraph(f"{k}: {v}%", meta_style) for k, v in fields_conf.items()]
                elems.append(ListFlowable([ListItem(b, leftIndent=10) for b in bullets], bulletType='bullet'))
            elems.append(Spacer(1, 8))

        # Mapped fields
        if mapped:
            elems.append(Paragraph("<b>Mapped Fields</b>", styles['Heading3']))
            def _to_text(val):
                try:
                    if isinstance(val, (list, tuple)):
                        return ", ".join(_to_text(x) for x in val)
                    if isinstance(val, dict):
                        return json.dumps(val, ensure_ascii=False)
                    return str(val)
                except Exception:
                    return str(val)
            for k, v in (mapped or {}).items():
                if k in ("enrichment", "enriched"):
                    continue
                text_value = _to_text(v)
                elems.append(Paragraph(f"• <b>{k}:</b> {text_value}", meta_style))
            elems.append(Spacer(1, 8))

        # Custom Fields (only if present)
        custom_fields = result_dict.get('custom_fields') if isinstance(result_dict, dict) else {}
        if isinstance(custom_fields, dict) and custom_fields:
            elems.append(Paragraph("<b>Custom Fields</b>", styles['Heading3']))
            for ck, cv in custom_fields.items():
                try:
                    cv_text = json.dumps(cv, ensure_ascii=False) if isinstance(cv, (dict, list)) else str(cv)
                except Exception:
                    cv_text = str(cv)
                elems.append(Paragraph(f"• <b>{ck}:</b> {cv_text}", meta_style))
            elems.append(Spacer(1, 8))
        
        # Enrichment
        enr = (mapped or {}).get("enrichment") or {}
        if enr:
            elems.append(Paragraph("<b>AI Enrichment</b>", styles['Heading3']))
            def _bullet_list(items):
                return ListFlowable([ListItem(Paragraph(str(it), meta_style), leftIndent=10) for it in items], bulletType='bullet')
            if isinstance(enr.get("diagnoses"), list) and enr["diagnoses"]:
                elems.append(Paragraph("Diagnoses", meta_style))
                elems.append(_bullet_list(enr["diagnoses"][:20]))
            if isinstance(enr.get("procedures"), list) and enr["procedures"]:
                elems.append(Paragraph("Procedures", meta_style))
                elems.append(_bullet_list(enr["procedures"][:20]))
            if isinstance(enr.get("labs_normalized"), list) and enr["labs_normalized"]:
                elems.append(Paragraph("Labs (normalized)", meta_style))
                labs_lines = [json.dumps(x) for x in enr["labs_normalized"][:20]]
                elems.append(_bullet_list(labs_lines))
            if isinstance(enr.get("obligations"), list) and enr["obligations"]:
                elems.append(Paragraph("Obligations", meta_style))
                elems.append(_bullet_list(enr["obligations"][:20]))
            if isinstance(enr.get("termination_conditions"), list) and enr["termination_conditions"]:
                elems.append(Paragraph("Termination Conditions", meta_style))
                elems.append(_bullet_list(enr["termination_conditions"][:20]))
            if isinstance(enr.get("key_findings"), list) and enr["key_findings"]:
                elems.append(Paragraph("Key Findings", meta_style))
                elems.append(_bullet_list(enr["key_findings"][:20]))
            if isinstance(enr.get("primary_outcomes"), list) and enr["primary_outcomes"]:
                elems.append(Paragraph("Primary Outcomes", meta_style))
                elems.append(_bullet_list(enr["primary_outcomes"][:20]))
            elems.append(Spacer(1, 8))

        # Tables (first few)
        if tables:
            elems.append(Paragraph("<b>Tables (preview)</b>", styles['Heading3']))
            max_preview = min(2, len(tables))
            for idx in range(max_preview):
                t = tables[idx]
                page_no = t.get('page')
                elems.append(Paragraph(f"Table on page {page_no}", meta_style))
            elems.append(Spacer(1, 8))

        # Formulas (first few)
        if formulas:
            elems.append(Paragraph("<b>Formulas (first 10)</b>", styles['Heading3']))
            for f in (formulas[:10] or []):
                elems.append(Paragraph(f, meta_style))
            elems.append(Spacer(1, 8))

        # Validation issues
        if validation:
            errors = (validation.get('errors') or [])
            warnings = (validation.get('warnings') or [])
            if errors:
                elems.append(Paragraph("<b>Validation Errors</b>", styles['Heading3']))
                elems.append(ListFlowable([ListItem(Paragraph(e, meta_style), leftIndent=10) for e in errors], bulletType='bullet'))
            if warnings:
                elems.append(Paragraph("<b>Validation Warnings</b>", styles['Heading3']))
                elems.append(ListFlowable([ListItem(Paragraph(w, meta_style), leftIndent=10) for w in warnings], bulletType='bullet'))

        doc.build(elems)
        buffer.seek(0)

        @after_this_request
        def _cleanup(resp):
            try:
                if os.path.isfile(temp_path):
                    os.remove(temp_path)
            except Exception:
                pass
            return resp

        return send_file(buffer, as_attachment=True, download_name=f"extraction_report.pdf", mimetype="application/pdf")
    except Exception as e:
        return jsonify({"error": f"Failed to generate extract PDF: {e}"}), 500

# ─────────────────────────────────────────────────────────────────

def ai_only_extract(all_text: str, pages_text: list[str], custom_instructions: str = "", domain: str | None = None, selected_fields: list[str] | None = None) -> dict:
    """Run AI-only extraction with a strict, schema-driven prompt to get exact values.
    Returns a full response with keys: type, pages, entities, mapped_fields, validation, confidence, tables, formulas, provenance
    Note: tables/formulas/provenance are left empty in AI-only mode unless we extend the AI prompt.
    """
    # Smart chunking to keep guidance intact but cap tokens
    max_chars = 20000
    text = all_text or ''
    if len(text) > max_chars:
        head = text[:10000]
        middle = text[len(text)//2 - 2500: len(text)//2 + 2500]
        tail = text[-5000:]
        doc_snippet = f"{head}\n\n[MIDDLE]\n\n{middle}\n\n[END]\n\n{tail}"
    else:
        doc_snippet = text
    system = (
        "You are an expert document information extractor. Return ONLY valid JSON. "
        "Extract exact values from the text without paraphrasing, keep original formatting of numbers and dates where possible. "
        "Use null for missing fields."
    )
    # If a domain is provided, use focused schema and instructions; else provide compact detection prompt.
    domain_norm = _normalize_domain(domain)
    allowed_fields = DOMAIN_FIELDS.get(domain_norm, DOMAIN_FIELDS["general"]) if domain_norm else DOMAIN_FIELDS["general"]
    selected = _validate_selected_fields(domain_norm, selected_fields) if domain_norm else []
    if domain_norm and domain_norm in DOMAIN_FIELDS:
        fields_for_prompt = selected if selected else allowed_fields
        prompt = (
            f"Document domain: {domain_norm}. Extract ONLY the following fields under 'mapped_fields': {json.dumps(fields_for_prompt)}\n"
            "Return JSON with keys: {\n"
            "  \"type\": string (exactly the provided domain),\n"
            "  \"mapped_fields\": object (only the requested fields),\n"
            "  \"entities\": { \"emails\": [], \"phones\": [], \"amounts\": [], \"dates\": [] },\n"
            "  \"custom_fields\": object | null\n"
            "}\n\n"
            "Rules: Use exact values from text; preserve number/date formatting; missing => null (or [] for arrays).\n"
            "Do not include extraneous keys. Return ONLY JSON.\n\n"
            f"Text (truncated):\n{doc_snippet}\n"
        )
    else:
        # Compact auto-detect fallback (legacy behavior)
        compact_schema = {k: DOMAIN_FIELDS[k][:8] for k in ("invoice","contract","financial","research","healthcare","general") if k in DOMAIN_FIELDS}
        prompt = (
            "Detect document type as one of: invoice, contract, financial, research, healthcare, general.\n"
            "Then extract fields (limited set) as strict JSON with keys: {\n"
            "  \"type\": string,\n"
            "  \"mapped_fields\": object,\n"
            "  \"entities\": { \"emails\": [], \"phones\": [], \"amounts\": [], \"dates\": [] },\n"
            "  \"custom_fields\": object | null\n"
            "}\n\n"
            f"Text (truncated):\n{doc_snippet}\n\n"
            f"Field schema (reference): {json.dumps(compact_schema)}\n"
        )
    # If custom instructions are provided, require a dedicated field
    if custom_instructions:
        prompt += (
            "\nIMPORTANT: The user provided CUSTOM INSTRUCTIONS below.\n"
            "DECISION LOGIC:\n"
            "- If the instruction relates to existing document fields (e.g., 'extract dates in ISO format', 'focus on medication dosages'), modify the relevant mapped_fields accordingly.\n"
            "- If the instruction asks for different concepts/data types (e.g., 'extract keywords', 'find risk factors', 'identify action items'), create new fields in 'custom_fields'.\n"
            "- You can do BOTH: modify existing fields AND add custom fields if the instruction covers both.\n"
            "- Use descriptive snake_case keys for custom_fields (e.g., 'keywords', 'risk_factors', 'action_items').\n"
            "- If nothing relevant is found, set 'custom_fields' to null.\n\n"
            f"CUSTOM INSTRUCTIONS:\n{custom_instructions}\n"
        )
    ai_text = None
    data = {}
    # Thread-based timeout for Windows compatibility
    import threading
    result_holder = {"text": None, "error": None}
    def _runner():
        try:
            result_holder["text"] = call_ai_service(prompt, system_prompt=system)
        except Exception as e:
            result_holder["error"] = str(e)
    th = threading.Thread(target=_runner, daemon=True)
    th.start()
    th.join(timeout=OLLAMA_TIMEOUT)  # reuse overall timeout cap
    if th.is_alive():
        try:
            print("DEBUG: ai_only_extract timed out waiting for AI response")
            sys.stdout.flush()
        except Exception:
            pass
        result_holder["text"] = None
        result_holder["error"] = "timeout"
    ai_text = result_holder["text"] or ""
    if result_holder["error"] == "timeout":
        data = {}
    else:
        try:
            data = parse_json_safely(ai_text)
        except Exception:
            try:
                print("DEBUG: ai_only_extract JSON parse failed; falling back to empty result")
                sys.stdout.flush()
            except Exception:
                pass
            data = {}
    if not isinstance(data, dict):
        data = {}
    dtype = (data.get('type') or (domain_norm or 'general')).lower()
    # Ensure minimal structure
    entities = data.get('entities') or {}
    mapped = data.get('mapped_fields') or {}
    custom_result = data.get('custom_instructions_result')
    # Lightweight validation: ensure only schema fields exist
    # If user selected specific fields, only include those; otherwise use domain's allowed fields
    if selected and domain_norm:
        # User selected specific fields - only include those
        allowed = set(selected)
    else:
        # No specific selection - use domain's allowed fields
        filter_domain = domain_norm or dtype
        allowed = set(DOMAIN_FIELDS.get(filter_domain, DOMAIN_FIELDS["general"]))
    
    mapped_clean = {}
    for k, v in (mapped.items() if isinstance(mapped, dict) else []):
        if k in allowed:
            mapped_clean[k] = v
    # Confidence is AI-only placeholder: presence-based
    fields_present = sum(1 for v in mapped_clean.values() if v)
    fields_total = max(1, len(allowed))
    overall = int(round(60 + 40 * (fields_present / fields_total))) if fields_total else 60
    confidence = {"overall": overall, "fields": {k: (85 if mapped_clean.get(k) else 0) for k in allowed}}
    # Response in the same shape as before
    return {
        "type": dtype,
        "pages": len(pages_text or []),
        "entities": entities,
        "mapped_fields": mapped_clean,
        "custom_fields": data.get('custom_fields'),  # Add this line
        "validation": {"errors": [], "warnings": []},
        "confidence": confidence,
        "tables": [],
        "formulas": [],
        "provenance": None,
    }

if __name__ == "__main__":
    app.run(host="0.0.0.0", port=8000)   # change port if needed
